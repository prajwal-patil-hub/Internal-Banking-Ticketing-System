"""Build the role-based SOP deck from the captured screenshots.

Layout rule from the prompt: one screen per slide, screenshot large enough to
read, numbered callouts on the elements that matter, and the same six fields on
every workflow slide (step, action, result, next).
"""
from __future__ import annotations

import json
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SHOTS = Path(sys.argv[1])
OUT = Path(sys.argv[2])


def _commit() -> str:
    """The commit the screenshots were taken from, read rather than typed.

    It was hand-written before and went stale the first time the app changed,
    which is exactly the claim the closing slide asks the reader to trust.
    """
    if len(sys.argv) > 3:
        return sys.argv[3]
    try:
        import subprocess
        return subprocess.run(["git", "rev-parse", "--short", "HEAD"],
                              capture_output=True, text=True, check=True,
                              cwd=Path(__file__).resolve().parent).stdout.strip() or "unknown"
    except Exception:
        return "unknown"


COMMIT = _commit()

# --- palette, taken from the application's own theme ------------------------
TEAL = RGBColor(0x0E, 0x4F, 0x4A)
TEAL_D = RGBColor(0x09, 0x38, 0x34)
CREAM = RGBColor(0xF3, 0xEC, 0xE0)
INK = RGBColor(0x1B, 0x2A, 0x33)
MUTE = RGBColor(0x5B, 0x6B, 0x75)
LINE = RGBColor(0xD6, 0xCE, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xC2, 0x5B, 0x2E)
OK = RGBColor(0x1E, 0x7A, 0x53)
WARN = RGBColor(0xB4, 0x7A, 0x14)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]

step_no = 0

# Callout coordinates come from the browser, not from guesses.
MANIFEST = json.loads((SHOTS / "manifest.json").read_text()) if (SHOTS / "manifest.json").exists() else {}

_FONT_CACHE: dict[tuple[float, bool], object] = {}


def _pil_font(pt: float, bold: bool):
    """Real font metrics, so wrapped text can be measured rather than guessed."""
    key = (round(pt, 1), bold)
    if key not in _FONT_CACHE:
        paths = (["/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"] if bold
                 else ["/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"])
        f = None
        for pth in paths:
            if Path(pth).exists():
                try:
                    f = ImageFont.truetype(pth, max(6, int(pt * 1.333)))
                except Exception:
                    f = None
        _FONT_CACHE[key] = f or ImageFont.load_default()
    return _FONT_CACHE[key]


_MEASURE_SS = 8   # supersample factor for text measurement


def wrapped_lines(text: str, width_in: float, pt: float, bold: bool = False) -> int:
    """How many lines `text` needs at `pt` inside `width_in` inches.

    Measured at 8x and scaled back. At 1x the font size is truncated to whole
    pixels — int(10.5 * 1.333) is 13px, ~7% under the true size — which made
    two-line legend labels report as one and let the next item overlap them.

    A 4% margin on the available width covers the difference between DejaVu
    (measured here) and Calibri (what PowerPoint will actually lay out). The
    error direction that matters is under-counting, since that overlaps text;
    over-counting only leaves a little extra space.
    """
    f = _pil_font(pt * _MEASURE_SS, bold)
    avail_px = width_in * 96.0 * _MEASURE_SS * 0.96
    d = ImageDraw.Draw(Image.new("RGB", (10, 10)))
    n = 0
    for para in text.split("\n"):
        n += 1
        line = ""
        for word in para.split():
            trial = (line + " " + word).strip()
            if d.textlength(trial, font=f) <= avail_px or not line:
                line = trial
            else:
                n += 1
                line = word
    return n


_TARGET_DPI = 200        # comfortably past what print or a screen resolves


def _fit_image(path: Path, display_w_emu: int):
    """Embed a screenshot at the size it is actually shown, not at capture size.

    Captures are 2880x1800 (a 1440x900 viewport at 2x) but they are displayed
    about 8.5in wide, so the full-resolution copy is roughly double what even
    print needs. Twenty-four of them made the file large enough that phone
    preview apps mis-composited it — showing one slide's picture on another,
    and leaving shapes behind between slides. Downscaling to 200dpi keeps the
    screenshots legible at 100% zoom and cuts the file several-fold.
    """
    import io
    target_px = int(display_w_emu / 914400 * _TARGET_DPI)
    im = Image.open(path)
    if im.width <= target_px:
        return str(path)
    im = im.convert("RGB").resize(
        (target_px, round(im.height * target_px / im.width)), Image.LANCZOS)
    # A UI screenshot is mostly flat fills and anti-aliased text, so a 256
    # colour palette is visually indistinguishable from truecolour here and
    # about a third of the size. Checked against the smallest type in the
    # captures before adopting.
    im = im.quantize(colors=256, dither=Image.FLOYDSTEINBERG)
    buf = io.BytesIO()
    im.save(buf, format="PNG", optimize=True)
    buf.seek(0)
    return buf


def marker(slide, cx, cy, n, diameter, fill, fontsize):
    """A numbered circle whose digit is actually centred.

    The default text insets are 0.1in each side. On a 0.29in circle that leaves
    0.09in for the glyph, so the number is squeezed out of position — which is
    exactly how the first version of this deck shipped.
    """
    sh = _rect(slide, cx - diameter / 2, cy - diameter / 2, diameter, diameter,
               fill, shape=MSO_SHAPE.OVAL)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Autofit off: with it on, a renderer may rescale the glyph and shift it
    # off centre. Line spacing pinned to 1.0 for the same reason — inherited
    # leading is added above the line, which pushes a vertically-centred
    # single character downward.
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.text = str(n)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    r = p.runs[0]
    r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(fontsize), True, WHITE, "Calibri"
    return sh


def _tb(slide, x, y, w, h, text, size=14, bold=False, color=INK,
        align=PP_ALIGN.LEFT, font="Calibri", spacing=1.15, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = font
    return box


def _rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    s.shadow.inherit = False
    return s


def _bg(slide, color=WHITE):
    """Paint the slide's own background rather than laying a shape over it.

    This used to add a full-bleed rectangle to every slide — 58 extra shapes
    sitting underneath everything else. A real background is what the format
    provides for this, and it gives weaker viewers one less large overlapping
    shape to composite.
    """
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def footer(slide, label):
    _rect(slide, 0, H - Inches(0.42), W, Inches(0.42), CREAM)
    _tb(slide, Inches(0.55), H - Inches(0.34), Inches(9), Inches(0.25),
        "SUCCESS Bank — Internal Ticketing · Standard Operating Procedure",
        size=9, color=MUTE)
    _tb(slide, W - Inches(2.4), H - Inches(0.34), Inches(1.85), Inches(0.25),
        label, size=9, color=MUTE, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- title -----
def title_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s, TEAL)
    _rect(s, 0, Inches(3.28), W, Inches(0.06), ACCENT)
    _tb(s, Inches(1.0), Inches(1.30), Inches(11), Inches(0.5),
        "STANDARD OPERATING PROCEDURE", size=15, bold=True, color=CREAM)
    # Broken across two lines on purpose. On one line it fits in Calibri and
    # nowhere else — Keynote and Google Slides substitute a wider face and the
    # title runs into the rule below it.
    _tb(s, Inches(1.0), Inches(1.70), Inches(11.6), Inches(1.45),
        "SUCCESS Bank\nInternal Ticketing", size=38, bold=True, color=WHITE, spacing=1.12)
    _tb(s, Inches(1.0), Inches(3.55), Inches(10), Inches(1.1),
        "A role-by-role operational guide.\n"
        "Every screen in this deck was captured from the running application.",
        size=17, color=CREAM, spacing=1.35)
    for i, (role, who) in enumerate([
        ("Branch User", "raises the ticket"),
        ("Agent", "works and resolves it"),
        ("Supervisor", "watches SLA and escalations"),
        ("Admin", "configures the system"),
        ("Auditor", "reads everything, changes nothing"),
    ]):
        x = Inches(1.0) + i * Inches(2.28)
        _rect(s, x, Inches(4.75), Inches(2.05), Inches(1.15), TEAL_D)
        _tb(s, x + Inches(0.18), Inches(4.95), Inches(1.7), Inches(0.3),
            role, size=12.5, bold=True, color=WHITE)
        _tb(s, x + Inches(0.18), Inches(5.28), Inches(1.72), Inches(0.55),
            who, size=9.5, color=CREAM, spacing=1.1)
    _tb(s, Inches(1.0), Inches(6.35), Inches(11), Inches(0.3),
        f"Version 1.0  \u00b7  Built at commit {COMMIT}", size=10, color=CREAM)


# -------------------------------------------------------------- section -----
def section(num, title, blurb):
    s = prs.slides.add_slide(BLANK)
    _bg(s, CREAM)
    _rect(s, 0, 0, Inches(0.28), H, TEAL)
    _tb(s, Inches(1.1), Inches(2.6), Inches(2), Inches(0.6),
        f"SECTION {num}", size=13, bold=True, color=ACCENT)
    _tb(s, Inches(1.1), Inches(3.05), Inches(10.5), Inches(0.9),
        title, size=34, bold=True, color=TEAL)
    _tb(s, Inches(1.1), Inches(4.15), Inches(9.5), Inches(1.0),
        blurb, size=15, color=MUTE, spacing=1.3)
    footer(s, f"Section {num}")


# --------------------------------------------------------------- content ----
def content(title, kicker=""):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    _rect(s, 0, 0, W, Inches(0.95), TEAL)
    _tb(s, Inches(0.55), Inches(0.2), Inches(11.5), Inches(0.35),
        title, size=23, bold=True, color=WHITE)
    if kicker:
        _tb(s, Inches(0.55), Inches(0.6), Inches(11.5), Inches(0.28),
            kicker, size=11.5, color=CREAM)
    return s


# ------------------------------------------------------- workflow slide -----
def workflow(shot, title, action, result, nxt, role_tag, extra_note=None):
    """One screen, one slide.

    Callout positions and labels come from `manifest.json`, which the capture
    script produced by asking the browser where each element actually is. The
    first version guessed percentages and put markers beside controls rather
    than on them.
    """
    global step_no
    step_no += 1

    # Callers write "{next}" instead of a hard-coded "Step 14". Hand-written
    # numbers went stale the moment a slide was inserted anywhere earlier —
    # three were already wrong before this deck was regenerated, all pointing
    # a reader at the slide after the one they meant.
    nxt = nxt.replace("{next}", f"Step {step_no + 1:02d}")
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    _rect(s, 0, 0, W, Inches(0.95), TEAL)
    _rect(s, 0, Inches(0.95), W, Inches(0.04), ACCENT)

    _tb(s, Inches(0.55), Inches(0.17), Inches(1.5), Inches(0.3),
        f"STEP {step_no:02d}", size=12, bold=True, color=ACCENT)
    _tb(s, Inches(0.55), Inches(0.46), Inches(9.2), Inches(0.4),
        title, size=20, bold=True, color=WHITE)
    _tb(s, W - Inches(3.3), Inches(0.35), Inches(2.75), Inches(0.3),
        role_tag.upper(), size=11, bold=True, color=CREAM, align=PP_ALIGN.RIGHT)

    entry = MANIFEST.get(shot.replace(".png", ""), {})
    callouts = entry.get("callouts", [])
    # A callout has to mark a *thing*. Twelve of them targeted the entire page
    # — one at 154%, a whole scrolling container — which put a numbered dot on
    # the sidebar with a label describing the screen as a whole. That is not an
    # annotation, it is decoration, and it was the bulk of the loose dots. Each
    # one's statement is already carried by the slide's Action/Result prose,
    # which is where a statement about the whole screen belongs.
    #
    # The threshold cleanly separates them: real targets here run under 16%
    # (the navigation menu is the largest), the discarded ones start at 78%.
    dropped = [c for c in callouts if c.get("w", 0) * c.get("h", 0) > 0.45]
    callouts = [c for c in callouts if c.get("w", 0) * c.get("h", 0) <= 0.45]
    for c in dropped:
        print(f"    dropped whole-screen callout on {shot}: {c['label'][:52]}")

    def _marker_y(c: dict) -> float:
        """Where the marker will actually be drawn, as a fraction of the image.

        A callout on a whole panel ('main', 'nav') is anchored near the top of
        the region — the centre of a full-height panel is empty space. The
        numbering has to sort on the same value, or a marker drawn at the top
        of the screen is numbered as though it sat in the middle. That is what
        put ② above ① on three slides.
        """
        return c.get("y0", 0.0) + 0.045 if c.get("h", 0) > 0.45 else c.get("y", 0.5)

    # Number in reading order, not in the order the capture script happened to
    # list them. y is banded first so two callouts at roughly the same height
    # run left-to-right instead of being ordered by a few stray pixels.
    # Rows are clustered, not banded. Any fixed grid has a knife edge: two
    # controls a third of an inch apart can straddle a boundary and count as
    # different rows, which numbered a button at the top-right before a search
    # box six inches to its left and barely lower. To a reader those are one
    # row, and the numbering ran backwards.
    #
    # Walking in y order and starting a new row only when the gap from the
    # row's first item exceeds ~0.5in has no boundary to straddle. Within a
    # row, left to right; the trailing y keeps a stacked column of form fields
    # in top-to-bottom order when they share an x.
    ROW_GAP = 0.10                      # fraction of image height, ~0.5in
    rows: list[list[dict]] = []
    for c in sorted(callouts, key=_marker_y):
        if rows and _marker_y(c) - _marker_y(rows[-1][0]) <= ROW_GAP:
            rows[-1].append(c)
        else:
            rows.append([c])
    callouts = [
        c
        for row in rows
        for c in sorted(row, key=lambda c: (c.get("x0", c.get("x", 0)), _marker_y(c)))
    ]

    img_x, img_y = Inches(0.45), Inches(1.3)
    path = SHOTS / shot
    # The image is fitted inside a box rather than pinned to one width, and its
    # height is read off the file rather than assumed to be 1440x900.
    #
    # The constant was correct for every screen captured at that viewport and
    # silently wrong for anything else. `add_picture` honours the real aspect,
    # so a taller screenshot drew past the height the markers were positioned
    # against — every marker landed high by a margin that grew down the slide.
    # Capping the height as well as the width keeps a tall screen inside the
    # slide and preserves the gap the note strip sits in.
    MAX_W, MAX_H = Inches(8.5), Inches(5.25)
    if path.exists():
        with Image.open(path) as _probe:
            ratio = _probe.height / _probe.width
    else:
        ratio = 900 / 1440
    img_w, img_h = MAX_W, Emu(int(MAX_W * ratio))
    if img_h > MAX_H:
        img_h = MAX_H
        img_w = Emu(int(MAX_H / ratio))
    if path.exists():
        _rect(s, img_x - Inches(0.025), img_y - Inches(0.025),
              img_w + Inches(0.05), img_h + Inches(0.05), LINE)
        s.shapes.add_picture(_fit_image(path, img_w), img_x, img_y, width=img_w)
    else:
        _rect(s, img_x, img_y, img_w, img_h, CREAM, LINE)
        _tb(s, img_x, img_y + img_h / 2, img_w, Inches(0.4),
            f"[ screen not captured: {shot} ]", size=13, color=ACCENT, align=PP_ALIGN.CENTER)

    # Markers sit just outside the element's left edge, vertically centred on
    # it. Centred on the element they cover the value the reader is meant to
    # read — the account number, the file name.
    D = Inches(0.30)
    for i, c in enumerate(callouts, start=1):
        # Same function the numbering sorted on, so the two cannot diverge.
        cy = img_y + Emu(int(img_h * _marker_y(c)))
        x0 = c.get("x0")
        if x0 is None:                      # older manifest: fall back to centre
            cx = img_x + Emu(int(img_w * c["x"]))
        else:
            left_edge = img_x + Emu(int(img_w * x0))
            cx = left_edge - D * 0.72       # outside, overlapping the border a little
            if cx < img_x + D * 0.5:        # no room on the left — sit inside it
                cx = img_x + Emu(int(img_w * x0)) + D * 0.62
        marker(s, cx, cy, i, D, ACCENT, 13)

    # ---- right column -----------------------------------------------------
    rx, rw = Inches(9.25), Inches(3.6)
    y = Inches(1.3)
    label_w = rw - Inches(0.36)

    if callouts:
        _tb(s, rx, y, rw, Inches(0.24), "ON THIS SCREEN", size=10, bold=True, color=ACCENT)
        y += Inches(0.32)
        for i, c in enumerate(callouts, start=1):
            lines = wrapped_lines(c["label"], label_w / 914400, 10.5)
            block_h = Inches(0.185) * lines
            # Same colour as the marker on the screenshot — the reader is
            # matching ① to ①, and a colour change breaks that thread.
            marker(s, rx + Inches(0.115), y + Inches(0.095), i, Inches(0.23), ACCENT, 10)
            _tb(s, rx + Inches(0.36), y, label_w, block_h,
                c["label"], size=10.5, color=INK, spacing=1.05)
            y += block_h + Inches(0.10)
        y += Inches(0.14)

    for lab, txt, col in (("ACTION", action, TEAL), ("RESULT", result, OK), ("NEXT", nxt, MUTE)):
        lines = wrapped_lines(txt, (rw - Inches(0.18)) / 914400, 11)
        body_h = Inches(0.175) * lines
        bar_h = Inches(0.2) + body_h
        _rect(s, rx, y, Inches(0.05), bar_h, col)
        _tb(s, rx + Inches(0.18), y, rw - Inches(0.18), Inches(0.18),
            lab, size=9.5, bold=True, color=col)
        _tb(s, rx + Inches(0.18), y + Inches(0.21), rw - Inches(0.18), body_h,
            txt, size=11, color=INK, spacing=1.08)
        y += bar_h + Inches(0.22)

    if extra_note:
        # Fills the gap between the screenshot and the footer exactly. It used
        # to be a fixed 0.5in box that started below the image and finished
        # underneath the footer bar, clipping the last line mid-word.
        gap_top = img_y + img_h + Inches(0.08)
        gap_h = (H - Inches(0.42)) - gap_top - Inches(0.04)
        _rect(s, img_x, gap_top, img_w, gap_h, CREAM)
        _tb(s, img_x + Inches(0.18), gap_top + Inches(0.07), img_w - Inches(0.36),
            gap_h - Inches(0.12), extra_note, size=10, color=INK,
            anchor=MSO_ANCHOR.MIDDLE)

    footer(s, f"Step {step_no:02d}")
    return s


def table(slide, x, y, w, headers, rows, col_w=None, head_fill=TEAL, size=10.5):
    n_r, n_c = len(rows) + 1, len(headers)
    gt = slide.shapes.add_table(n_r, n_c, x, y, w, Inches(0.4) * n_r).table
    if col_w:
        for i, cw in enumerate(col_w):
            gt.columns[i].width = cw
    for c, htxt in enumerate(headers):
        cell = gt.cell(0, c)
        cell.text = htxt
        cell.fill.solid()
        cell.fill.fore_color.rgb = head_fill
        p = cell.text_frame.paragraphs[0]
        p.runs[0].font.size, p.runs[0].font.bold = Pt(size), True
        p.runs[0].font.color.rgb = WHITE
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else CREAM
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.size = Pt(size - 0.5)
            p.runs[0].font.color.rgb = INK
    return gt


def chip(slide, x, y, w, h, text, fill, tcol=WHITE, size=11, bold=True):
    """A rounded label box.

    Every paragraph must be styled, not just the first. `tf.text = "a\\nb"`
    splits into two paragraphs, and styling only paragraphs[0] leaves the
    second line at the 18pt black default — which overflowed every two-line
    chip on the swimlane slide.
    """
    sh = _rect(slide, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.1
        for r in p.runs:
            r.font.size, r.font.bold, r.font.color.rgb = Pt(size), bold, tcol
    return sh


def arrow(slide, x, y, w, h=Inches(0.16), color=MUTE):
    a = _rect(slide, x, y, w, h, color, shape=MSO_SHAPE.RIGHT_ARROW)
    return a


def arrow_left(slide, x, y, w, h=Inches(0.16), color=MUTE):
    return _rect(slide, x, y, w, h, color, shape=MSO_SHAPE.LEFT_ARROW)


def arrow_up(slide, x, y, w, h=Inches(0.22), color=MUTE):
    return _rect(slide, x, y, w, h, color, shape=MSO_SHAPE.UP_ARROW)


# ===========================================================================
# 01 — Title and overview
# ===========================================================================
title_slide()

s = content("How to use this deck", "Read it in order the first time; use Section 14 as a desk reference after that")
_tb(s, Inches(0.6), Inches(1.4), Inches(5.6), Inches(0.3), "WHAT THIS IS", size=11, bold=True, color=ACCENT)
_tb(s, Inches(0.6), Inches(1.75), Inches(5.6), Inches(2.2),
    "A step-by-step operational guide. Each role is taken from signing in to a "
    "finished ticket, using screens captured from the running application.\n\n"
    "Every workflow slide follows the same shape: what you do, what the system "
    "does back, and who picks it up next.", size=13, color=INK, spacing=1.3)
_tb(s, Inches(0.6), Inches(4.0), Inches(5.6), Inches(0.3), "WHAT IT IS NOT", size=11, bold=True, color=ACCENT)
_tb(s, Inches(0.6), Inches(4.35), Inches(5.6), Inches(1.6),
    "Not a feature tour, and not a specification. Where the system does not do "
    "something people often expect — approvals, in-app notifications, ticket "
    "merging — this deck says so rather than staying silent.",
    size=13, color=INK, spacing=1.3)
_rect(s, Inches(6.9), Inches(1.4), Inches(5.85), Inches(4.6), CREAM)
_tb(s, Inches(7.25), Inches(1.7), Inches(5.2), Inches(0.3), "FIND YOUR ROLE", size=11, bold=True, color=ACCENT)
for i, (r, sec, line) in enumerate([
    ("Branch User", "Section 6", "You raise tickets and follow them"),
    ("Agent", "Section 7", "You work and resolve them"),
    ("Supervisor", "Section 8", "You watch SLA and escalations"),
    ("Admin", "Section 9", "You configure the system"),
    ("Auditor", "Section 10", "You review; you change nothing"),
]):
    y = Inches(2.15) + i * Inches(0.72)
    _rect(s, Inches(7.25), y, Inches(0.05), Inches(0.55), TEAL)
    _tb(s, Inches(7.45), y, Inches(2.2), Inches(0.28), r, size=13, bold=True, color=TEAL)
    _tb(s, Inches(7.45), y + Inches(0.26), Inches(4.6), Inches(0.28), line, size=10.5, color=MUTE)
    _tb(s, Inches(11.3), y + Inches(0.02), Inches(1.2), Inches(0.28), sec, size=10.5, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
footer(s, "Overview")

# ===========================================================================
# 02 — System overview
# ===========================================================================
section("02", "What the system is for", "One record for every operational problem, from the moment it is raised to the moment it is closed.")

s = content("The problem it replaces", "Why a shared record beats a shared inbox")
for i, (t, b, col) in enumerate([
    ("Without it",
     "A blocked card is reported by phone. A duplicated debit goes to somebody's "
     "inbox. Nobody can say how long either has been open, who owns it, or what "
     "was decided.", ACCENT),
    ("With it",
     "Every problem gets a number, an owner, a deadline and a history. Months "
     "later an auditor can answer who changed what, and when.", OK),
]):
    x = Inches(0.6) + i * Inches(6.3)
    _rect(s, x, Inches(1.4), Inches(5.9), Inches(2.1), CREAM)
    _rect(s, x, Inches(1.4), Inches(0.06), Inches(2.1), col)
    _tb(s, x + Inches(0.3), Inches(1.65), Inches(5.3), Inches(0.3), t, size=15, bold=True, color=col)
    _tb(s, x + Inches(0.3), Inches(2.05), Inches(5.3), Inches(1.3), b, size=12.5, color=INK, spacing=1.3)

_tb(s, Inches(0.6), Inches(3.85), Inches(12), Inches(0.3), "HOW A TICKET ARRIVES", size=11, bold=True, color=ACCENT)
for i, (t, d) in enumerate([
    ("Portal", "A branch user fills in the form and attaches evidence"),
    ("Email", "A message to the support mailbox becomes a ticket automatically"),
]):
    x = Inches(0.6) + i * Inches(6.3)
    _rect(s, x, Inches(4.25), Inches(5.9), Inches(1.0), WHITE, LINE)
    _tb(s, x + Inches(0.3), Inches(4.45), Inches(2), Inches(0.3), t, size=13, bold=True, color=TEAL)
    _tb(s, x + Inches(0.3), Inches(4.75), Inches(5.3), Inches(0.4), d, size=11, color=MUTE)
_tb(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.9),
    "Email intake is implemented but off by default (IMAP_ENABLED=false) and needs a real mailbox. "
    "Everything else in this deck works out of the box.", size=11, color=MUTE, spacing=1.25)
footer(s, "Overview")

# ===========================================================================
# 03 — Roles and permissions
# ===========================================================================
section("03", "Roles and permissions", "Five roles, one per user. This is the whole model — there are no per-user overrides.")

s = content("Who can do what", "Enforced server-side in core/authz.py — the interface only hides what the API would refuse anyway")
table(s, Inches(0.5), Inches(1.3), Inches(12.3),
      ["Role", "Tickets", "Assignment", "Users & org", "Audit", "Escalations"],
      [["Branch User", "Own only — raise, comment, attach", "—", "—", "—", "—"],
       ["Agent", "Any in scope — progress, resolve", "Assign to a person", "—", "—", "—"],
       ["Supervisor", "As agent", "Also auto-assign, rules, leave", "Read directory", "—", "Yes"],
       ["Admin", "As agent", "As supervisor, plus the delay", "Full control", "Yes", "Yes"],
       ["Auditor", "Read only — no writes at all", "—", "—", "Yes", "—"]],
      col_w=[Inches(1.75), Inches(3.35), Inches(2.85), Inches(1.95), Inches(1.05), Inches(1.35)])
_rect(s, Inches(0.5), Inches(4.25), Inches(12.3), Inches(1.15), CREAM)
_rect(s, Inches(0.5), Inches(4.25), Inches(0.06), Inches(1.15), ACCENT)
_tb(s, Inches(0.8), Inches(4.45), Inches(11.8), Inches(0.8),
    "Super admin is a second tier on top of Admin, not a sixth role. Only a super admin can create "
    "another super admin or change one's password — otherwise any admin could take over the account.",
    size=12, color=INK, spacing=1.25)
_tb(s, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.8),
    "The auditor is read-only everywhere. It can open any ticket, dashboard and report, and every "
    "attempt to write is refused with a 403 — including comments and attachments.",
    size=12, color=MUTE, spacing=1.25)
footer(s, "Roles")

# ===========================================================================
# 04 — Lifecycle
# ===========================================================================
section("04", "The ticket lifecycle", "Nine statuses. The API enforces which moves are legal — an illegal jump is refused, not merely discouraged.")

s = content("Status flow", "Reproduced from VALID_TRANSITIONS in ticket_service.py — these are the only legal moves")

# --- main line -------------------------------------------------------------
main = [("new", TEAL), ("acknowledged", TEAL), ("assigned", TEAL),
        ("in_progress", TEAL), ("resolved", OK), ("closed", MUTE)]
cw, gap = Inches(1.72), Inches(0.32)
x0, ymain = Inches(0.62), Inches(1.5)
xs = []
for i, (label, col) in enumerate(main):
    x = x0 + i * (cw + gap)
    xs.append(x)
    chip(s, x, ymain, cw, Inches(0.6), label, col, size=11)
    if i < len(main) - 1:
        arrow(s, x + cw + Inches(0.03), ymain + Inches(0.22), gap - Inches(0.06))

# --- branch states: both hang off in_progress and both return to it --------
ip_cx = xs[3] + cw / 2
brow = Inches(2.62)
_rect(s, ip_cx - Inches(0.02), ymain + Inches(0.6), Inches(0.04), Inches(0.22), LINE)

pair = [("on_hold", WARN, ip_cx - cw - Inches(0.16)),
        ("escalated", ACCENT, ip_cx + Inches(0.16))]
left_cx = pair[0][2] + cw / 2
right_cx = pair[1][2] + cw / 2
_rect(s, left_cx, ymain + Inches(0.78), right_cx - left_cx, Inches(0.04), LINE)
for label, col, bx in pair:
    _rect(s, bx + cw / 2 - Inches(0.02), ymain + Inches(0.78), Inches(0.04),
          brow - (ymain + Inches(0.78)), LINE)
    chip(s, bx, brow, cw, Inches(0.52), label, col, size=10.5)

_tb(s, Inches(0.62), brow + Inches(0.08), Inches(3.4), Inches(0.5),
    "Both pause the ticket and\nboth return to in_progress.", size=9.5, color=MUTE, spacing=1.15)

# --- reopened: closed goes back to assigned, via its own row ---------------
rrow = Inches(3.72)
rh = Inches(0.52)
rmid = rrow + rh / 2
closed_cx = xs[5] + cw / 2
assigned_cx = xs[2] + cw / 2
chip(s, xs[2], rrow, cw, rh, "reopened", WARN, size=10.5)

# closed drops to the reopened row and runs *left* into reopened. The arrowhead
# has to point at reopened; a right-pointing head here reads reopened→closed,
# which is the opposite of what the state machine allows.
_rect(s, closed_cx - Inches(0.02), ymain + Inches(0.6), Inches(0.04),
      rmid - (ymain + Inches(0.6)), LINE)
bar_x0 = xs[2] + cw + Inches(0.30)
_rect(s, bar_x0, rmid - Inches(0.02), closed_cx - bar_x0 + Inches(0.02),
      Inches(0.04), LINE)
arrow_left(s, xs[2] + cw + Inches(0.04), rmid - Inches(0.08), Inches(0.26),
           Inches(0.16), LINE)

# reopened rises back into assigned, in its own column so nothing is crossed
_rect(s, assigned_cx - Inches(0.02), ymain + Inches(0.82), Inches(0.04),
      rrow - (ymain + Inches(0.82)), LINE)
arrow_up(s, assigned_cx - Inches(0.09), ymain + Inches(0.60), Inches(0.18),
         Inches(0.24), LINE)
_tb(s, Inches(0.62), rrow + Inches(0.06), Inches(3.4), Inches(0.5),
    "closed can be reopened; it\nreturns to assigned.", size=9.5, color=MUTE, spacing=1.15)

# --- the two rules ---------------------------------------------------------
_tb(s, Inches(0.62), Inches(4.6), Inches(12.1), Inches(0.26),
    "TWO RULES PEOPLE GET WRONG", size=10.5, bold=True, color=ACCENT)
for i, (ttl, d) in enumerate([
    ("closed is reachable from any open state",
     "Closing early is a withdrawal — the problem went away, or it was raised in error."),
    ("resolved is NOT reachable from assigned",
     "The ticket must have been worked first, so 'resolved' means something."),
]):
    y = Inches(4.95) + i * Inches(0.78)
    _rect(s, Inches(0.62), y, Inches(12.1), Inches(0.66), CREAM)
    _tb(s, Inches(0.9), y + Inches(0.09), Inches(11.6), Inches(0.24), ttl, size=12, bold=True, color=TEAL)
    _tb(s, Inches(0.9), y + Inches(0.35), Inches(11.6), Inches(0.26), d, size=10.5, color=MUTE)
footer(s, "Lifecycle")

# ===========================================================================
# 05 — Swimlane
# ===========================================================================
section("05", "End to end, across the roles", "The seams matter more than the steps. This is where a ticket changes hands.")

s = content("Who holds the ticket, and when", "Left to right in time; each band is one role")
lanes = [("BRANCH USER", TEAL), ("AGENT", TEAL_D), ("SUPERVISOR", ACCENT), ("SYSTEM", MUTE)]
lane_h, lane_y0 = Inches(1.12), Inches(1.35)
for i, (name, col) in enumerate(lanes):
    y = lane_y0 + i * lane_h
    _rect(s, Inches(0.5), y, Inches(1.85), lane_h - Inches(0.06), col)
    _tb(s, Inches(0.62), y + Inches(0.36), Inches(1.65), Inches(0.3), name, size=10.5, bold=True, color=WHITE)
    _rect(s, Inches(2.4), y, Inches(10.4), lane_h - Inches(0.06), CREAM if i % 2 == 0 else WHITE, LINE)
# Nine steps, not eight. The old diagram had the system auto-assigning at
# creation, which is no longer true and was the single most misleading thing
# in the deck: it showed a machine making the decision a supervisor makes.
steps = [
    (0, 0, "Raise ticket\n+ attach evidence"),
    (3, 1, "Number it,\nstamp SLA"),
    (2, 2, "Assign, or\nauto-assign"),
    (1, 3, "Take it,\ninvestigate"),
    (1, 4, "Reply\n+ attach fix"),
    (3, 5, "SLA breach?\nescalate"),
    (2, 6, "Review,\nreassign"),
    (1, 7, "Resolve"),
    (0, 8, "See resolution,\nclose or reopen"),
]
# Nine columns across the same 10.4in band: the pitch shrinks rather than the
# band growing, so the lanes still line up with the role labels.
sw_w, sw_pitch, sw_h = Inches(0.98), Inches(1.145), Inches(0.76)
sw_x0 = Inches(2.52)


def _sw_box(lane, col_i):
    x = sw_x0 + col_i * sw_pitch
    y = lane_y0 + lane * lane_h + Inches(0.16)
    return x, y, x + sw_w, y + sw_h


# Connectors first, so the chips sit on top of them rather than under.
for (la, ca, _), (lb, cb, _) in zip(steps, steps[1:]):
    _, ya, xa1, ya1 = _sw_box(la, ca)
    xb, yb, _, yb1 = _sw_box(lb, cb)
    mid_a, mid_b = (ya + ya1) / 2, (yb + yb1) / 2
    gx = xa1 + (xb - xa1) / 2          # centre of the gutter between the two
    _rect(s, xa1, mid_a - Inches(0.015), gx - xa1, Inches(0.03), LINE)
    top, bot = min(mid_a, mid_b), max(mid_a, mid_b)
    _rect(s, gx - Inches(0.015), top, Inches(0.03), bot - top, LINE)
    _rect(s, gx, mid_b - Inches(0.015), xb - gx - Inches(0.12), Inches(0.03), LINE)
    arrow(s, xb - Inches(0.14), mid_b - Inches(0.06), Inches(0.14), Inches(0.12), LINE)

for lane, col_i, label in steps:
    x, y, _, _ = _sw_box(lane, col_i)
    chip(s, x, y, sw_w, sw_h, label, lanes[lane][1], size=8)

_tb(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9),
    "The connectors are the handoffs. Two of them are a person deciding: the supervisor choosing an "
    "owner, and the branch user closing or reopening. The rest happen because the system moved the "
    "ticket, or because the next role was already watching. Note what the system does not do — it "
    "numbers the ticket and starts the clock, but it does not choose who works it.",
    size=11.5, color=MUTE, spacing=1.3)
footer(s, "End to end")

# ===========================================================================
# 06 — Branch user
# ===========================================================================
section("06", "Branch user workflow", "You raise the problem, supply the evidence, answer questions, and see the resolution.")

workflow("00-login.png", "Sign in",  "Enter your bank email and password, then select Sign in.",
   "You land on your dashboard. If your account has a second factor enabled, "
   "you are asked for a 6-digit code first.",
   "{next} — read your dashboard.", "Branch User")

workflow("10-branch-dashboard.png", "Read your dashboard",  "Check the tiles for anything of yours that is breaching or still open.",
   "Each tile opens the exact list it counts — a tile reading 3 opens three tickets.",
   "{next} — open the ticket list.", "Branch User")

workflow("11-branch-tickets.png", "See only your own tickets",  "Scan the list, or filter to what you are looking for.",
   "You see only tickets you raised. This is enforced by the server, not hidden by the page.",
   "{next} — raise a new ticket.", "Branch User")

workflow("12-branch-create-empty.png", "Open the new ticket form",  "Select New Ticket from the ticket list.",
   "An empty form opens. Nothing is submitted until you choose Create Ticket.",
   "{next} — fill it in and attach evidence.", "Branch User")

workflow("13-branch-create-filled.png", "Describe it, and attach the evidence",  "Complete the form and attach a screenshot, statement or spreadsheet. "
   "Up to 15 MB per file; images, PDF, Word, Excel, text and CSV.",
   "Files are held in the browser and uploaded the moment the ticket is created — "
   "so evidence arrives with the report, not after it.",
   "{next} — submit and follow it.", "Branch User")

workflow("15-branch-ticket-detail.png", "Follow your ticket",  "Open the ticket from your list to see where it stands.",
   "The ticket was numbered, given SLA deadlines and assigned automatically — "
   "no one had to triage it by hand.",
   "{next} — answer questions and read the resolution.", "Branch User")

workflow("16-branch-ticket-comments.png", "Reply, and read the resolution",  "Answer any question the agent asks, attaching more evidence if needed.",
   "When the agent resolves it, their explanation and any file they attached "
   "appear here together.",
   "If it is fixed, the ticket is closed. If not, it can be reopened.", "Branch User")

s = content("What a branch user cannot see", "Three rules worth knowing before you ask where something went")
for i, (t, d, col) in enumerate([
    ("Internal notes — and their attachments",
     "Agents can mark a note internal. You will not see the note, and you will not see any "
     "file attached to it. It is not hidden in the page; the server refuses to send it.", ACCENT),
    ("Other people's tickets",
     "You see only tickets you raised. Opening someone else's by its link returns nothing.", TEAL),
    ("Team-wide numbers",
     "SLA health, category and department breakdowns and AI metrics are for agents and above. "
     "Your dashboard shows your own three counts instead.", TEAL),
]):
    y = Inches(1.35) + i * Inches(1.42)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(1.24), CREAM)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(1.24), col)
    _tb(s, Inches(0.95), y + Inches(0.18), Inches(11.5), Inches(0.32), t, size=15, bold=True, color=col)
    _tb(s, Inches(0.95), y + Inches(0.58), Inches(11.5), Inches(0.75), d, size=12.5, color=INK, spacing=1.3)
_tb(s, Inches(0.6), Inches(5.85), Inches(12.1), Inches(1.2),
    "If you believe a ticket of yours is missing, it is far more likely that it was raised by a "
    "colleague than that it was deleted — tickets are never deleted, and every change is recorded "
    "in the audit trail.", size=12.5, color=MUTE, spacing=1.3)
footer(s, "Branch User")

# ===========================================================================
# 07 — Agent
# ===========================================================================
section("07", "Agent workflow", "You pick the ticket up, investigate, keep the requester informed, and resolve it.")

workflow("20-agent-dashboard.png", "Start from the dashboard",  "Sign in and read the KPI strip before opening anything.",
   "Every tile is a live count and opens the exact list behind it.",
   "{next} — open the queue.", "Agent")

workflow("21-agent-tickets.png", "Work the queue",  "Filter to unassigned or to your own, and choose what to work on.",
   "Agents see every ticket in their org scope, not just their own.",
   "{next} — start with what is breaching.", "Agent")

workflow("22-agent-breached.png", "Deal with breaches first",  "Select the SLA Breached tile on the dashboard.",
   "The list is filtered to exactly the tickets the tile counted — the number "
   "on the card and the length of this list always agree.",
   "{next} — open one and investigate.", "Agent")

workflow("23-agent-ticket-detail.png", "Investigate the ticket",  "Read the description and the attached evidence, then move the ticket to In Progress.",
   "The status change is recorded in the audit trail with your name, and the "
   "first-response clock stops.",
   "{next} — reply, and attach the fix.", "Agent")

workflow("24-agent-ticket-comments.png", "Reply — and attach the fix to your reply",  "Write what you found and attach the corrected statement or screenshot.",
   "Files attached here belong to this reply, so the requester sees your fix "
   "beside the answer that explains it.",
   "{next} — resolve, or escalate.", "Agent")

s = content("Resolve, or escalate — and what each means", "Both are one action on the ticket page")
for i, (t, d, col) in enumerate([
    ("Resolve", "Use when the problem is fixed. Only available once the ticket has been worked — "
                "you cannot jump straight from Assigned to Resolved, which is what makes 'resolved' "
                "mean something.", OK),
    ("Escalate", "Use when it needs someone else. The ticket moves to Escalated, is reassigned to the "
                 "least-loaded holder of the target role, and an escalation event is recorded. "
                 "Escalating twice for the same reason does nothing — the second is suppressed.", ACCENT),
    ("Put on hold", "Use when you are waiting on someone outside the team. The ticket stays open and "
                    "keeps counting toward your queue.", WARN),
]):
    y = Inches(1.4) + i * Inches(1.55)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(1.32), CREAM, LINE)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(1.32), col)
    _tb(s, Inches(0.95), y + Inches(0.2), Inches(11.4), Inches(0.3), t, size=15, bold=True, color=col)
    _tb(s, Inches(0.95), y + Inches(0.58), Inches(11.4), Inches(0.68), d, size=12, color=INK, spacing=1.28)
_tb(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.6),
    "Closing is always available, from any open state — it means the problem went away or the ticket "
    "was raised in error, not that it was solved.", size=11.5, color=MUTE, spacing=1.25)
footer(s, "Agent")

# ===========================================================================
# 08 — Supervisor
# ===========================================================================
section("08", "Supervisor workflow", "You watch the deadlines and the escalations, and step in before a breach becomes a complaint.")

workflow("30-supervisor-dashboard.png", "Watch the team's position",  "Read the strip; anything breaching needs an owner today.",
   "Supervisors see the same tiles as agents, plus the SLA monitor and escalation queue in the menu.",
   "{next} — open the SLA monitor.", "Supervisor")

workflow("31-supervisor-sla.png", "Read the SLA monitor",  "Work down from the tickets nearest their deadline.",
   "At risk means due within the hour. Breached means the deadline has already passed "
   "and, if a rule matched, escalation has already fired.",
   "{next} — review escalations.", "Supervisor")

workflow("32-supervisor-escalations.png", "Review the escalation queue",  "Check what escalated and whether the target is acting on it.",
   "Escalations arrive here whether raised by hand or fired automatically by the "
   "SLA worker — both run the same engine, so the evidence is identical.",
   "Reassign if the target is wrong; otherwise the agent resolves it.", "Supervisor")

workflow("33-supervisor-unassigned.png", "Decide who works it",
   "Open a ticket with no owner and select Assign.",
   "Raising a ticket stamps its SLA deadlines but does not choose an owner. That is "
   "deliberate: who carries the work is a shift decision, and a machine making it at "
   "2am is how tickets end up with whoever happens to be idlest rather than whoever "
   "should have them.",
   "{next} — pick a person, or let the router pick.", "Supervisor")

workflow("34-supervisor-assign-list.png", "Pick a person, or let the router pick",
   "Choose a name, or select Auto-assign to take the lightest queue.",
   "Each name shows the number of open tickets that person is already carrying — the "
   "same number the router ranks on, so you can see what you are overriding. Anyone on "
   "leave is labelled with their return date and sorted to the bottom. You can still "
   "assign to them knowingly; auto-assign will not.",
   "The agent picks it up from their queue.", "Supervisor",
   extra_note="Auto-assign is supervisor and above; an agent may assign a specific ticket.")

s = content("What the router actually does", "Three steps, and it stops at the first that yields somebody")
_tb(s, Inches(0.62), Inches(1.3), Inches(12.1), Inches(0.5),
    "Auto-assign is not a black box. It runs the same search every time, and each step "
    "falls through to the next rather than failing.", size=12.5, color=INK, spacing=1.25)
for i, (t, d) in enumerate([
    ("1 — A category rule",
     "If a rule names an owner for this category, they get it. Optional: most categories have none."),
    ("2 — Someone in the ticket's branch",
     "Of the people left, whoever in that branch is carrying the fewest open tickets."),
    ("3 — Anyone assignable",
     "Failing both, the lightest open queue anywhere."),
]):
    y = Inches(2.0) + i * Inches(0.92)
    _rect(s, Inches(0.62), y, Inches(12.1), Inches(0.8), CREAM)
    _tb(s, Inches(0.95), y + Inches(0.12), Inches(3.6), Inches(0.3), t, size=12.5, bold=True, color=TEAL)
    _tb(s, Inches(4.6), y + Inches(0.12), Inches(7.9), Inches(0.6), d, size=11.5, color=INK, spacing=1.2)

_tb(s, Inches(0.62), Inches(4.9), Inches(12.1), Inches(0.28),
    "TWO RULES THAT HOLD AT EVERY STEP", size=10.5, bold=True, color=ACCENT)
for i, (ttl, d) in enumerate([
    ("Agents before supervisors",
     "Ranking on workload alone sends everything to whoever is idlest, which is reliably a "
     "supervisor — they carry no queue of their own — and frontline work would skip the agents entirely."),
    ("Nobody on leave, ever",
     "A category rule naming someone who is away is skipped rather than obeyed. Parking tickets "
     "on an absent person is worse than having no rule."),
]):
    y = Inches(5.22) + i * Inches(0.86)
    _rect(s, Inches(0.62), y, Inches(12.1), Inches(0.74), CREAM)
    _tb(s, Inches(0.9), y + Inches(0.08), Inches(11.6), Inches(0.24), ttl, size=12, bold=True, color=TEAL)
    _tb(s, Inches(0.9), y + Inches(0.34), Inches(11.6), Inches(0.36), d, size=10.5, color=MUTE, spacing=1.15)
footer(s, "Supervisor")

s = content("The ticket nobody picked up", "Why moving assignment to a person does not risk the SLA")
_rect(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(1.3), CREAM)
_rect(s, Inches(0.6), Inches(1.35), Inches(0.06), Inches(1.3), ACCENT)
_tb(s, Inches(0.95), Inches(1.58), Inches(11.4), Inches(0.9),
    "SLA deadlines are stamped when the ticket is raised, and the clock runs from that "
    "moment. If assignment waited for a person and nobody was on shift, a ticket raised "
    "overnight could breach — and then escalate — without ever having had an owner.",
    size=13, color=INK, spacing=1.3)
for i, (t, d) in enumerate([
    ("The safety net",
     "A background worker assigns anything still unassigned after a set delay, using the same router."),
    ("You set the delay",
     "An admin sets it while the system is running — two hours by default, anywhere from fifteen minutes to a week."),
    ("It never overrides you",
     "It only touches tickets with no owner. It never reassigns one somebody has already given out."),
]):
    y = Inches(3.0) + i * Inches(0.92)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(0.8), WHITE, LINE)
    _tb(s, Inches(0.9), y + Inches(0.12), Inches(3.0), Inches(0.3), t, size=12.5, bold=True, color=TEAL)
    _tb(s, Inches(3.9), y + Inches(0.12), Inches(8.6), Inches(0.6), d, size=11.5, color=INK, spacing=1.2)
_tb(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.6),
    "So the window is yours to triage in, not a gap. Shorten it if tickets sit too long; "
    "lengthen it if the system is assigning work your supervisors wanted to place themselves.",
    size=11.5, color=MUTE, spacing=1.25)
footer(s, "Supervisor")

s = content("How escalation actually fires", "Nobody has to be watching for this to happen")
_rect(s, Inches(0.6), Inches(1.35), Inches(12.1), Inches(1.15), CREAM)
_tb(s, Inches(0.95), Inches(1.6), Inches(11.4), Inches(0.7),
    "Every 5 minutes a background worker looks for tickets past their resolution deadline. "
    "For each one it finds, it marks the breach and evaluates the escalation rules.",
    size=13, color=INK, spacing=1.3)
for i, (t, d) in enumerate([
    ("1 — Rule matched", "The most specific rule wins. A category-specific rule beats a catch-all, and priority_threshold is read as a minimum, so a 'high' rule also covers critical."),
    ("2 — Ticket moved", "Status becomes Escalated and it is reassigned to the least-loaded holder of the target role."),
    ("3 — Evidence written", "An escalation_events row records the trigger, the rule and the target."),
    ("4 — People told", "The target and the manager list are emailed after the change commits — a failed email never rolls back the escalation."),
]):
    y = Inches(2.75) + i * Inches(0.9)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(0.78), WHITE, LINE)
    _tb(s, Inches(0.9), y + Inches(0.11), Inches(3.0), Inches(0.28), t, size=12.5, bold=True, color=TEAL)
    _tb(s, Inches(3.9), y + Inches(0.11), Inches(8.6), Inches(0.6), d, size=11, color=INK, spacing=1.2)
_tb(s, Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.5),
    "A ticket will not escalate twice for the same trigger — otherwise the worker would raise a new "
    "event every five minutes until somebody resolved it.", size=11.5, color=MUTE, spacing=1.25)
footer(s, "Supervisor")

# ===========================================================================
# 09 — Admin
# ===========================================================================
section("09", "Admin workflow", "You decide who exists, what they may do, and how the organisation is shaped.")

workflow("41-admin-users.png", "Manage users",  "Create a user, set their role, and place them in an org unit.",
   "The role decides what they may do. There are no per-user permission overrides.",
   "{next} — shape the organisation.", "Admin")

workflow("46-admin-users-availability.png", "See who is available",
   "Read the Availability column before wondering why work is not reaching someone.",
   "Availability and Active are different columns because they mean different things. "
   "Active is whether the account works at all; availability is whether the router "
   "sends new tickets there. Someone on leave can still sign in and finish what they "
   "already hold.",
   "Record a leave window with the Leave button.", "Admin")

workflow("47-admin-leave-dialog.png", "Record a leave window",
   "Enter the first and last day away, then save.",
   "Auto-assign skips them for that window and starts including them again the day "
   "after it ends — nobody has to remember to switch them back on. A supervisor can "
   "still assign to them deliberately, which is why they stay in the list rather than "
   "disappearing from it.",
   "{next} — shape the organisation.", "Admin",
   extra_note="Supervisors can set leave too — rota changes should not wait for an admin.")

workflow("42-admin-org.png", "Shape the organisation",  "Define hierarchy levels, then units within them.",
   "The org tree drives ticket visibility: a user sees their unit's subtree, "
   "plus anything assigned to them personally.",
   "{next} — maintain the branch network.", "Admin")

workflow("43-admin-branches.png", "Maintain the branch network",  "Add branches, set managers and capacity, and mark degraded ones.",
   "Ticket counts are computed per request, never stored — a counter that "
   "drifts is wrong forever with nothing to reveal it.",
   "{next} — pull reports.", "Admin")

workflow("44-admin-reports.png", "Pull reports",  "Choose a period, then export.",
   "The export is generated from what is on screen, so it matches the filters "
   "you applied rather than silently re-running an unfiltered query.",
   "{next} — protect your own account.", "Admin")

workflow("45-admin-security.png", "Turn on your second factor",  "Open Security, scan the QR code, and enter one code to confirm.",
   "MFA is only switched on once a correct code proves the app is working — "
   "so a failed scan cannot lock you out. Ten single-use recovery codes are "
   "shown once, and never again.",
   "Save the recovery codes somewhere safe before leaving this screen.", "Admin")

s = content("Save the recovery codes", "This is the one screen you cannot go back to")
_rect(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.6), CREAM)
_rect(s, Inches(0.6), Inches(1.4), Inches(0.06), Inches(1.6), ACCENT)
_tb(s, Inches(0.95), Inches(1.68), Inches(11.4), Inches(1.15),
    "Ten recovery codes are shown once, when you switch MFA on. Only their hashes are stored, so "
    "nobody — including an administrator — can show them to you again. Each works once, in place of "
    "a code from your app.", size=14, color=INK, spacing=1.35)
for i, (t, d) in enumerate([
    ("Lost your phone, still have codes", "Use any unused recovery code at sign-in. It works once, then it is spent."),
    ("Lost both", "An administrator clears your enrolment, and you set it up again from scratch."),
    ("Lost the only super admin", "There is no route through the interface, by design. The flag has to be set on the database row."),
]):
    y = Inches(3.3) + i * Inches(0.95)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(0.8), WHITE, LINE)
    _tb(s, Inches(0.95), y + Inches(0.12), Inches(4.3), Inches(0.28), t, size=12.5, bold=True, color=TEAL)
    _tb(s, Inches(5.4), y + Inches(0.12), Inches(7.0), Inches(0.6), d, size=11.5, color=INK, spacing=1.2)
footer(s, "Admin")

# ===========================================================================
# 10 — Auditor
# ===========================================================================
section("10", "Auditor workflow", "You can open everything and change nothing. Every write is refused by the server.")

workflow("50-auditor-dashboard.png", "See the whole picture",  "Sign in and read the dashboard.",
   "Auditors have full visibility. What they do not have is any way to change "
   "what they are looking at.",
   "{next} — open the audit trail.", "Auditor")

workflow("51-auditor-audit-log.png", "Read the audit trail",  "Filter to the entity or person you are reviewing.",
   "Every state change writes a row: actor, role, IP, request id, and the "
   "values before and after.",
   "{next} — inspect any ticket.", "Auditor")

workflow("52-auditor-tickets.png", "Inspect any ticket",  "Open any ticket and read its comments, attachments and timeline.",
   "There is no scope limit — an auditor sees every ticket, not a subset. "
   "Internal notes are visible to them too. Write controls are not offered, "
   "and would be refused if called directly.",
   "This is the end of the auditor's workflow — there is nothing to submit.", "Auditor")

s = content("A caveat worth recording", "Honest limitation of the audit trail as it stands today")
_rect(s, Inches(0.6), Inches(1.5), Inches(12.1), Inches(2.0), CREAM)
_rect(s, Inches(0.6), Inches(1.5), Inches(0.06), Inches(2.0), ACCENT)
_tb(s, Inches(0.95), Inches(1.8), Inches(11.4), Inches(0.35),
    "The audit trail is append-only by convention, not by enforcement.", size=16, bold=True, color=ACCENT)
_tb(s, Inches(0.95), Inches(2.3), Inches(11.4), Inches(1.1),
    "The application only ever inserts rows. But there is no database trigger and no revoked "
    "permission, so anything holding the database credentials could alter history. Closing this "
    "is the first item on the outstanding work list.", size=13, color=INK, spacing=1.35)
_tb(s, Inches(0.6), Inches(3.9), Inches(12.1), Inches(0.9),
    "It is recorded here rather than left out, because an auditor is exactly the person who needs to "
    "know how much weight the trail can carry.", size=12.5, color=MUTE, spacing=1.3)
footer(s, "Auditor")

# ===========================================================================
# 11 — Attachments
# ===========================================================================
section("11", "Attachments", "Where files can be added, who can see them, and the one rule people assume wrong.")

s = content("Three places a file can be attached", "The flow is the same in each: choose the files, then they upload once the thing they belong to exists")
for i, (t, d, who) in enumerate([
    ("With the ticket, when it is raised",
     "Evidence travels with the report rather than arriving afterwards.", "Branch user, agent"),
    ("With a reply",
     "The file belongs to that reply, so a fix sits beside the answer explaining it.", "Anyone who can comment"),
    ("From the attachments panel",
     "The file index for the whole ticket; files from replies are marked as such.", "Anyone who can write"),
]):
    y = Inches(1.4) + i * Inches(1.28)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(1.08), CREAM, LINE)
    _tb(s, Inches(0.95), y + Inches(0.16), Inches(6.4), Inches(0.3), t, size=13.5, bold=True, color=TEAL)
    _tb(s, Inches(0.95), y + Inches(0.52), Inches(6.4), Inches(0.45), d, size=11.5, color=INK, spacing=1.2)
    _tb(s, Inches(7.7), y + Inches(0.16), Inches(4.7), Inches(0.3), "WHO", size=9.5, bold=True, color=ACCENT)
    _tb(s, Inches(7.7), y + Inches(0.44), Inches(4.7), Inches(0.4), who, size=11.5, color=INK)
_tb(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.3),
    "Limits: 15 MB per file. Images, PDF, Word, Excel, text and CSV are accepted. Executables and "
    "archives are refused outright — there is no malware scanner, and an unscanned archive is the "
    "usual way something nasty arrives.", size=12.5, color=MUTE, spacing=1.3)
footer(s, "Attachments")

s = content("The rule people get wrong", "An internal note's attachment is as invisible as the note")
_rect(s, Inches(0.6), Inches(1.4), Inches(5.85), Inches(2.4), CREAM)
_rect(s, Inches(0.6), Inches(1.4), Inches(0.06), Inches(2.4), TEAL)
_tb(s, Inches(0.95), Inches(1.65), Inches(5.2), Inches(0.32), "What agents assume", size=14, bold=True, color=TEAL)
_tb(s, Inches(0.95), Inches(2.08), Inches(5.2), Inches(1.5),
    "\"The note is internal, so the requester will not read my text — but the file I attached "
    "is just a file on the ticket.\"", size=12.5, color=INK, spacing=1.3)
_rect(s, Inches(6.9), Inches(1.4), Inches(5.85), Inches(2.4), CREAM)
_rect(s, Inches(6.9), Inches(1.4), Inches(0.06), Inches(2.4), OK)
_tb(s, Inches(7.25), Inches(1.65), Inches(5.2), Inches(0.32), "What actually happens", size=14, bold=True, color=OK)
_tb(s, Inches(7.25), Inches(2.08), Inches(5.2), Inches(1.5),
    "The file is hidden too. It is filtered from the list and refused on download — "
    "and the refusal is a 404, so the response does not even confirm the note exists.",
    size=12.5, color=INK, spacing=1.3)
_tb(s, Inches(0.6), Inches(4.15), Inches(12.1), Inches(1.2),
    "This matters because the file is usually the sensitive part. Hiding the note while serving its "
    "attachment would leak exactly what the flag exists to withhold.", size=13, color=INK, spacing=1.3)
_rect(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.0), WHITE, LINE)
_tb(s, Inches(0.95), Inches(5.5), Inches(11.4), Inches(0.65),
    "Related: you cannot attach a file to somebody else's comment. Without that rule an agent could "
    "hang a file off the requester's own message, where it would read as something they had sent.",
    size=11.5, color=MUTE, spacing=1.25)
footer(s, "Attachments")

# ===========================================================================
# 12 — Assistant
# ===========================================================================
section("12", "The AI assistant", "It answers questions about the work in front of you. It cannot change anything.")

s = content("What it can and cannot do", "Three properties are enforced by the server, not by the prompt")
for i, (t, d, col) in enumerate([
    ("It cannot see more than you can",
     "Context is fetched through the same visibility rules as the rest of the application. "
     "Point it at a ticket your role cannot open and it is told the ticket is unavailable — "
     "no title, no number, nothing.", OK),
    ("It says when it does not know",
     "Rather than producing generic advice. An assistant that invents ticket facts in a bank "
     "is worse than one that declines.", OK),
    ("It cannot act",
     "It has no write access. It will not change a status, assign a ticket or post a comment. "
     "Everything it produces is for you to act on.", ACCENT),
]):
    y = Inches(1.4) + i * Inches(1.55)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(1.32), CREAM, LINE)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(1.32), col)
    _tb(s, Inches(0.95), y + Inches(0.18), Inches(11.4), Inches(0.3), t, size=14, bold=True, color=col)
    _tb(s, Inches(0.95), y + Inches(0.56), Inches(11.4), Inches(0.7), d, size=12, color=INK, spacing=1.25)
_tb(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.7),
    "If it reports that it cannot connect, the cause is almost always the local model rather than the "
    "application. The message it shows includes the specific fix.", size=11.5, color=MUTE, spacing=1.25)
footer(s, "Assistant")


# ===========================================================================
# 13 — Knowledge base
# ===========================================================================
section("13", "The knowledge base",
        "Ask a question, get an answer with the passage it came from. "
        "Administrators decide which documents exist and who may search them.")

s = content("Why this is not just a search box", "Three properties the server enforces, not the prompt")
for i, (t, d, col) in enumerate([
    ("It can only quote what you may already read",
     "Access is applied to the search itself, before anything is generated. A passage your role "
     "has not been granted is never retrieved, so it cannot appear in an answer — you get the same "
     "response as if the document did not exist.", OK),
    ("Every claim carries its source",
     "Each factual sentence is marked with the passage it came from, and you can open that "
     "document, section and page to check it. If the assistant produces a source that does not "
     "exist, the sentence it supported is removed before you see it.", OK),
    ("It refuses rather than guessing",
     "When the documents do not cover your question it says so and stops. That is the correct "
     "answer, not a failure — and it names what would help, such as asking an administrator "
     "whether the relevant document has been uploaded.", ACCENT),
]):
    y = Inches(1.4) + i * Inches(1.55)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(1.32), CREAM, LINE)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(1.32), col)
    _tb(s, Inches(0.95), y + Inches(0.18), Inches(11.4), Inches(0.3), t, size=14, bold=True, color=col)
    _tb(s, Inches(0.95), y + Inches(0.56), Inches(11.4), Inches(0.7), d, size=12, color=INK, spacing=1.25)
_tb(s, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.7),
    "Find it in the sidebar under Knowledge Base. Agents, supervisors and administrators can ask "
    "questions; only administrators can upload documents or change who may search a collection.",
    size=11.5, color=MUTE, spacing=1.25)
footer(s, "Knowledge base")

workflow("70-kb-answer.png", "Ask the knowledge base",
   "Type your question in plain words and select Ask. There is no query syntax and no filters.",
   "You get an answer with a confidence band, and every factual sentence marked with the passage "
   "it came from. Cited sources are listed below with their document, section and page.",
   "{next} — read what it retrieved but did not use.", "Agent & above")

workflow("70-kb-answer.png", "Check the answer",
   "Read the cited sources, and the passages listed under Also retrieved, not cited.",
   "The second list is what the search considered and the answer did not use. Seeing both tells "
   "you how much was weighed; a thin list is a reason to treat the answer carefully.",
   "{next} — what happens when it has nothing.", "Agent & above",
   "The confidence band is computed from the evidence, not asked of the assistant.")

workflow("72-kb-abstain.png", "When it declines",
   "Ask something the uploaded documents do not cover.",
   "It reports that it has no grounded answer and explains why, rather than producing something "
   "plausible. Try different wording, or ask an administrator whether the document exists.",
   "{next} — administrators: stocking the base.", "Agent & above")

workflow("71-kb-documents.png", "Upload and grant access",
   "Select a collection, choose Upload document, then set which roles may search it.",
   "The file is parsed, split into passages and indexed. It becomes searchable only once every "
   "passage is indexed — a document is never half-available.",
   "{next} — the two states worth knowing.", "Admin only")

s = content("Two states you will meet", "Both are the system being honest rather than broken")
for i, (t, d, col) in enumerate([
    ("\u201cNo roles granted \u2014 not searchable\u201d",
     "A new collection is readable by nobody until you grant a role. It will accept uploads and "
     "answer nothing until you do, so the list says so rather than leaving you to wonder why "
     "search returns silence.", ACCENT),
    ("A document marked Failed, with the reason",
     "Most often a scanned PDF with no text layer, or the local model being unreachable. The "
     "previous version of that document keeps answering throughout \u2014 select Re-index once the "
     "cause is fixed.", ACCENT),
]):
    y = Inches(1.5) + i * Inches(1.9)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(1.62), CREAM, LINE)
    _rect(s, Inches(0.6), y, Inches(0.06), Inches(1.62), col)
    _tb(s, Inches(0.95), y + Inches(0.2), Inches(11.4), Inches(0.3), t, size=14, bold=True, color=col)
    _tb(s, Inches(0.95), y + Inches(0.6), Inches(11.4), Inches(0.9), d, size=12, color=INK, spacing=1.25)
_tb(s, Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.1),
    "Branch users and auditors cannot search the knowledge base at all, and there is no navigation "
    "entry for them. It holds internal staff procedure, and an auditor is an oversight role rather "
    "than a working one.", size=11.5, color=MUTE, spacing=1.3)
footer(s, "Knowledge base")

# ===========================================================================
# 14 — Scenarios
# ===========================================================================
section("14", "Common scenarios", "Four situations, start to finish, naming who does what at each hop.")

for title, kicker, rows in [
    ("Scenario A — a duplicated debit",
     "The everyday case: raised with evidence, worked, resolved",
     [("Branch user", "Raises the ticket with a screenshot, the statement and a CSV of the transactions"),
      ("System", "Numbers it and stamps both SLA deadlines. It does not pick an owner"),
      ("Supervisor", "Assigns it — by hand, or with Auto-assign to take the lightest queue"),
      ("Agent", "Moves it to In Progress, checks the evidence, confirms the duplicate"),
      ("Agent", "Replies with the corrected statement attached to that reply, then resolves"),
      ("Branch user", "Sees the explanation and the corrected file together, and the ticket closes")]),
    ("Scenario B — nobody picked it up",
     "What happens when a deadline passes and no one is watching",
     [("System", "The SLA worker finds the ticket past its resolution deadline"),
      ("System", "Marks the breach, matches the most specific escalation rule"),
      ("System", "Moves it to Escalated, reassigns to the least-loaded holder of the target role"),
      ("System", "Writes an escalation event, then emails the target and the managers"),
      ("Supervisor", "Sees it in the escalation queue and confirms someone is on it")]),
    ("Scenario C — the agent needs more information",
     "The two-way exchange that used to happen over email",
     [("Agent", "Replies asking for the account number and a clearer screenshot"),
      ("Branch user", "Answers on the ticket and attaches the new screenshot"),
      ("Agent", "Continues, adding an internal note with findings the requester should not see"),
      ("Branch user", "Sees neither the internal note nor the file attached to it"),
      ("Agent", "Resolves with a public explanation")]),
    ("Scenario D — it was not actually fixed",
     "Reopening, and what it does to the counts",
     [("Branch user", "Reads the resolution and finds the problem still occurring"),
      ("Branch user", "Reopens the ticket, explaining what is still wrong"),
      ("System", "The ticket returns to the open set and counts against the dashboard again"),
      ("Agent", "Picks it back up from Assigned and works it a second time"),
      ("Agent", "Resolves again — the whole history stays on the one ticket")]),
]:
    s = content(title, kicker)
    # The row pitch is derived from how many rows there are, not fixed. Adding
    # the supervisor step to Scenario A made it six rows, and at the old fixed
    # 1.02in pitch the last one ran underneath the footer bar.
    top, bottom = Inches(1.32), H - Inches(0.62)
    pitch = min(Inches(1.02), (bottom - top) // max(1, len(rows)))
    row_h = pitch - Inches(0.16)
    for i, (who, what) in enumerate(rows):
        y = top + i * pitch
        col = TEAL if who not in ("System",) else MUTE
        _rect(s, Inches(0.6), y, Inches(2.25), row_h, col)
        _tb(s, Inches(0.78), y, Inches(2.0), row_h, who, size=12, bold=True, color=WHITE,
            anchor=MSO_ANCHOR.MIDDLE)
        _rect(s, Inches(2.95), y, Inches(9.78), row_h, CREAM, LINE)
        _tb(s, Inches(3.25), y, Inches(9.2), row_h, what, size=12, color=INK, spacing=1.2,
            anchor=MSO_ANCHOR.MIDDLE)
    footer(s, "Scenarios")

# ===========================================================================
# 15 — Quick reference
# ===========================================================================
section("15", "Quick reference", "The desk copy. Statuses, who does what, and where things live.")

s = content("Status reference", "Who can move a ticket into each state, and what it means")
table(s, Inches(0.5), Inches(1.25), Inches(12.3),
      ["Status", "Means", "Set by"],
      [["new", "Raised, not yet looked at", "System, on creation"],
       ["acknowledged", "Seen, not yet owned", "Agent"],
       ["assigned", "Has an owner", "Supervisor, or agent; system after the delay"],
       ["in_progress", "Actively being worked", "Agent"],
       ["on_hold", "Waiting on someone outside the team", "Agent"],
       ["escalated", "Handed up — by hand or by SLA breach", "Agent, or the SLA worker"],
       ["resolved", "Fixed; only from in_progress or escalated", "Agent"],
       ["closed", "Finished, or withdrawn; reachable from any open state", "Agent, or the requester"],
       ["reopened", "Was not actually fixed; returns to the open set", "Requester or agent"]],
      col_w=[Inches(2.1), Inches(6.6), Inches(3.6)], size=11)
footer(s, "Reference")

s = content("Where to click", "The shortest path to the thing you need")
table(s, Inches(0.5), Inches(1.25), Inches(12.3),
      ["I want to…", "Go to", "Who can"],
      [["Raise a problem", "Tickets → New Ticket", "Everyone"],
       ["Follow my ticket", "Tickets → open it", "Everyone (own only, for branch users)"],
       ["Find what is breaching", "Dashboard → SLA Breached tile", "Agent and above"],
       ["Give a ticket an owner", "Ticket → Assignee → Assign", "Agent and above"],
       ["Let the router choose", "Ticket → Assignee → Auto-assign", "Supervisor, admin"],
       ["Say someone is away", "Users → Leave", "Supervisor, admin"],
       ["See what escalated", "Escalations", "Supervisor, admin"],
       ["Check deadlines", "SLA Monitor", "Supervisor, admin"],
       ["Add or change a user", "Users", "Admin"],
       ["Change the org tree", "Org Management", "Admin"],
       ["Export figures", "Reports → CSV / PDF / Excel", "Supervisor, admin, auditor"],
       ["Review who changed what", "Audit", "Admin, auditor"],
       ["Turn on my second factor", "Security", "Everyone"]],
      col_w=[Inches(4.0), Inches(4.3), Inches(4.0)], size=11)
footer(s, "Reference")

s = content("What this system does not do", "Stated plainly, so nobody waits for a feature that is not coming today")
for i, (t, d) in enumerate([
    ("No approval workflow", "Nothing routes for sign-off before proceeding."),
    ("No in-app notifications", "Notifications are email only. There is no bell icon."),
    ("No ticket merging", "Duplicates can be marked as such, but two tickets cannot be combined into one."),
    ("No customer portal", "Every account is an employee account; branch staff raise tickets on customers' behalf."),
    ("No per-user permissions", "One role per user. Anything finer is expressed through the org tree."),
]):
    y = Inches(1.35) + i * Inches(1.0)
    _rect(s, Inches(0.6), y, Inches(12.1), Inches(0.84), CREAM, LINE)
    _tb(s, Inches(0.95), y + Inches(0.14), Inches(4.0), Inches(0.3), t, size=12.5, bold=True, color=ACCENT)
    _tb(s, Inches(5.1), y + Inches(0.14), Inches(7.4), Inches(0.6), d, size=11.5, color=INK, spacing=1.2)
_tb(s, Inches(0.6), Inches(6.5), Inches(12.1), Inches(0.4),
    "If one of these would help, it belongs on the ticket list — not in a workaround.",
    size=11.5, color=MUTE)
footer(s, "Reference")

# ---- close -----------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
_bg(s, TEAL)
_tb(s, Inches(1.0), Inches(2.6), Inches(11.6), Inches(1.05),
    "Questions, or something here that does not\nmatch what you see?",
    size=26, bold=True, color=WHITE, spacing=1.2)
_tb(s, Inches(1.0), Inches(3.85), Inches(11), Inches(1.4),
    f"This deck describes the application at commit {COMMIT}. If a screen has changed, the deck is "
    "wrong and not the application — raise it so this can be recaptured.\n\n"
    "Operational procedures — deploying, restoring, what to do at 3am — are in docs/runbook.md.",
    size=15, color=CREAM, spacing=1.4)

prs.save(str(OUT))
print(f"saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
