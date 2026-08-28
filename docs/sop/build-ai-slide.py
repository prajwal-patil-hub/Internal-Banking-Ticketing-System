"""Build the single 'how to ask the AI assistant' slide as a standalone deck.

Same visual language as the SOP deck so it can be pasted straight into it,
but self-contained: one slide, one screenshot, its own file.

    python docs/sop/build-ai-slide.py <shots-dir> <out.pptx>
"""
from __future__ import annotations

import io
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SHOTS = Path(sys.argv[1])
OUT = Path(sys.argv[2])

TEAL = RGBColor(0x0E, 0x4F, 0x4A)
CREAM = RGBColor(0xF3, 0xEC, 0xE0)
INK = RGBColor(0x1B, 0x2A, 0x33)
MUTE = RGBColor(0x5B, 0x6B, 0x75)
LINE = RGBColor(0xD6, 0xCE, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xC2, 0x5B, 0x2E)
OK = RGBColor(0x1E, 0x7A, 0x53)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background.fill
bg.solid()
bg.fore_color.rgb = WHITE


def rect(x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def tb(x, y, w, h, text, size=12, bold=False, color=INK,
       align=PP_ALIGN.LEFT, spacing=1.2, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        for r in p.runs:
            r.font.size, r.font.bold, r.font.color.rgb, r.font.name = (
                Pt(size), bold, color, "Calibri")
    return box


def marker(cx, cy, n, d, fill, fontsize):
    """Numbered circle with the digit actually centred — zero insets, no autofit."""
    sh = rect(cx - d / 2, cy - d / 2, d, d, fill, shape=MSO_SHAPE.OVAL)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.text = str(n)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    r = p.runs[0]
    r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(fontsize), True, WHITE, "Calibri"


# ---- header ----------------------------------------------------------------
rect(0, 0, W, Inches(0.95), TEAL)
rect(0, Inches(0.95), W, Inches(0.04), ACCENT)
tb(Inches(0.55), Inches(0.17), Inches(4), Inches(0.3),
   "THE AI ASSISTANT", size=12, bold=True, color=ACCENT)
tb(Inches(0.55), Inches(0.46), Inches(9.2), Inches(0.4),
   "Asking it a question", size=20, bold=True, color=WHITE)
tb(W - Inches(3.3), Inches(0.35), Inches(2.75), Inches(0.3),
   "EVERY ROLE", size=11, bold=True, color=CREAM, align=PP_ALIGN.RIGHT)

# ---- screenshot ------------------------------------------------------------
entry = json.loads((SHOTS / "manifest.json").read_text())["60-ai-assistant"]
img_x, img_y, img_w = Inches(0.45), Inches(1.30), Inches(8.50)
src = SHOTS / entry["file"]
im = Image.open(src)
target = int(img_w / 914400 * 200)          # 200dpi is past what print resolves
im = im.convert("RGB").resize((target, round(im.height * target / im.width)), Image.LANCZOS)
im = im.quantize(colors=256, dither=Image.FLOYDSTEINBERG)
buf = io.BytesIO()
im.save(buf, format="PNG", optimize=True)
buf.seek(0)
img_h = Inches(round(img_w / 914400 * 900 / 1440, 4))
rect(img_x - Inches(0.025), img_y - Inches(0.025),
     img_w + Inches(0.05), img_h + Inches(0.05), LINE)
s.shapes.add_picture(buf, img_x, img_y, width=img_w)


def marker_y(c):
    return c.get("y0", 0.0) + 0.045 if c.get("h", 0) > 0.45 else c.get("y", 0.5)


# Reading order: rows clustered by a real gap, then left to right.
rows: list[list[dict]] = []
for c in sorted(entry["callouts"], key=marker_y):
    if rows and marker_y(c) - marker_y(rows[-1][0]) <= 0.10:
        rows[-1].append(c)
    else:
        rows.append([c])
callouts = [c for row in rows for c in sorted(row, key=lambda c: c.get("x0", 0))]

D = Inches(0.30)
for i, c in enumerate(callouts, start=1):
    cy = img_y + Emu(int(img_h * marker_y(c)))
    left = img_x + Emu(int(img_w * c["x0"]))
    cx = left - D * 0.72
    if cx < img_x + D * 0.5:
        cx = left + D * 0.62
    marker(cx, cy, i, D, ACCENT, 13)

# ---- right column ----------------------------------------------------------
rx, rw = Inches(9.25), Inches(3.6)
y = Inches(1.30)

tb(rx, y, rw, Inches(0.24), "ON THIS SCREEN", size=10, bold=True, color=ACCENT)
y += Inches(0.32)
for i, c in enumerate(callouts, start=1):
    lines = 2 if len(c["label"]) > 34 else 1
    h = Inches(0.185) * lines
    marker(rx + Inches(0.115), y + Inches(0.095), i, Inches(0.23), ACCENT, 10)
    tb(rx + Inches(0.36), y, rw - Inches(0.36), h, c["label"], size=10.5, color=INK, spacing=1.05)
    y += h + Inches(0.10)

y += Inches(0.16)
rect(rx, y, Inches(0.05), Inches(1.62), TEAL)
tb(rx + Inches(0.18), y, rw - Inches(0.18), Inches(0.18),
   "ASK IT LIKE THIS", size=9.5, bold=True, color=TEAL)
tb(rx + Inches(0.18), y + Inches(0.23), rw - Inches(0.18), Inches(1.35),
   "“Summarise this ticket for the customer”\n"
   "“Which category has the most load right now?”\n"
   "“What is breaching today?”\n"
   "“What are the SLA policies for critical tickets?”\n"
   "“Suggest next steps on this one”",
   size=10.5, color=INK, spacing=1.25)
y += Inches(1.80)

rect(rx, y, Inches(0.05), Inches(0.92), ACCENT)
tb(rx + Inches(0.18), y, rw - Inches(0.18), Inches(0.18),
   "IT WILL NOT", size=9.5, bold=True, color=ACCENT)
tb(rx + Inches(0.18), y + Inches(0.23), rw - Inches(0.18), Inches(0.66),
   "Change a status, assign a ticket or post a comment. Everything it "
   "produces is for you to act on.",
   size=10.5, color=INK, spacing=1.15)

# ---- note under the screenshot ---------------------------------------------
gap_top = img_y + img_h + Inches(0.08)
gap_h = (H - Inches(0.42)) - gap_top - Inches(0.04)
rect(img_x, gap_top, img_w, gap_h, CREAM)
tb(img_x + Inches(0.18), gap_top + Inches(0.07), img_w - Inches(0.36), gap_h - Inches(0.12),
   "It answers only from what your role can already open — ask about a hidden "
   "ticket and it says so.",
   size=10, color=INK, anchor=MSO_ANCHOR.MIDDLE)

# ---- footer ----------------------------------------------------------------
rect(0, H - Inches(0.42), W, Inches(0.42), CREAM)
tb(Inches(0.55), H - Inches(0.34), Inches(9), Inches(0.25),
   "SUCCESS Bank — Internal Ticketing · Standard Operating Procedure",
   size=9, color=MUTE)
tb(W - Inches(2.4), H - Inches(0.34), Inches(1.85), Inches(0.25),
   "AI Assistant", size=9, color=MUTE, align=PP_ALIGN.RIGHT)

prs.save(str(OUT))
print(f"saved {OUT} (1 slide, {len(callouts)} callouts)")
