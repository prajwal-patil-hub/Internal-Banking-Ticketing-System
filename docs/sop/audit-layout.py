"""Check a .pptx for layout faults, so review is not 58 rounds of eyeballing.

Reads the shape tree and reports, per slide:

  BLEED     a shape extends past the slide edge
  OVERFLOW  text needs more height than its box has, so it spills
  OVERLAP   two text-bearing boxes cover the same area
  TINY      a numbered marker whose glyph cannot fit its circle

Text height is measured with PIL against DejaVu, the same way the builder
sizes its blocks. PowerPoint lays out in Calibri, which is narrower, so this
over-estimates slightly — findings are conservative in the safe direction.
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

SRC = Path(sys.argv[1])
EMU = 914400.0
SS = 8
_D = ImageDraw.Draw(Image.new("RGB", (10, 10)))
_CACHE: dict = {}


def font(pt: float, bold: bool):
    key = (round(pt, 1), bold)
    if key not in _CACHE:
        p = ("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold
             else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
        _CACHE[key] = ImageFont.truetype(p, max(6, int(pt * 1.333)))
    return _CACHE[key]


def text_lines(text: str, width_in: float, pt: float, bold: bool) -> int:
    f = font(pt * SS, bold)
    avail = width_in * 96.0 * SS * 0.96
    n = 0
    for para in text.split("\n"):
        n += 1
        line = ""
        for w in para.split():
            t = (line + " " + w).strip()
            if _D.textlength(t, font=f) <= avail or not line:
                line = t
            else:
                n += 1
                line = w
    return n


def widest(text: str, pt: float, bold: bool) -> float:
    """Widest unbreakable word, in inches."""
    f = font(pt * SS, bold)
    words = [w for para in text.split("\n") for w in para.split()] or [""]
    return max(_D.textlength(w, font=f) for w in words) / (96.0 * SS)


prs = Presentation(str(SRC))
SW, SH = prs.slide_width / EMU, prs.slide_height / EMU
findings: list[str] = []

for idx, slide in enumerate(prs.slides, start=1):
    boxes = []           # (x, y, w, h, text, name) for overlap checking
    for sh in slide.shapes:
        if sh.left is None or sh.width is None:
            continue
        x, y = sh.left / EMU, sh.top / EMU
        w, h = sh.width / EMU, sh.height / EMU

        if x < -0.01 or y < -0.01 or x + w > SW + 0.01 or y + h > SH + 0.01:
            findings.append(
                f"s{idx:02d} BLEED    {sh.shape_type} at ({x:.2f},{y:.2f}) "
                f"{w:.2f}x{h:.2f} vs slide {SW:.2f}x{SH:.2f}")

        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        txt = tf.text.strip()
        if not txt:
            continue

        ml = (tf.margin_left if tf.margin_left is not None else Emu(91440)) / EMU
        mr = (tf.margin_right if tf.margin_right is not None else Emu(91440)) / EMU
        mt = (tf.margin_top if tf.margin_top is not None else Emu(45720)) / EMU
        mb = (tf.margin_bottom if tf.margin_bottom is not None else Emu(45720)) / EMU
        inner_w = max(0.01, w - ml - mr)
        inner_h = max(0.01, h - mt - mb)

        need = 0.0
        maxw = 0.0
        for p in tf.paragraphs:
            if not p.runs:
                continue
            r0 = p.runs[0]
            pt = r0.font.size.pt if r0.font.size else 18.0
            bold = bool(r0.font.bold)
            ptxt = "".join(r.text for r in p.runs)
            if not ptxt.strip():
                continue
            sp = p.line_spacing if isinstance(p.line_spacing, float) else 1.15
            need += text_lines(ptxt, inner_w, pt, bold) * (pt * sp) / 72.0
            maxw = max(maxw, widest(ptxt, pt, bold))

        # A single word wider than the box cannot wrap and will stick out.
        if maxw > inner_w + 0.02:
            findings.append(
                f"s{idx:02d} OVERFLOW word {maxw:.2f}in > box {inner_w:.2f}in "
                f"— {txt[:44]!r}")
        elif need > inner_h + 0.06:
            findings.append(
                f"s{idx:02d} OVERFLOW needs {need:.2f}in in {inner_h:.2f}in "
                f"— {txt[:44]!r}")

        # A marker is a small circle holding a number; check the glyph fits.
        if len(txt) <= 2 and txt.isdigit() and w < 0.5:
            r0 = tf.paragraphs[0].runs[0]
            pt = r0.font.size.pt if r0.font.size else 18.0
            gw = widest(txt, pt, bool(r0.font.bold))
            if gw > inner_w:
                findings.append(
                    f"s{idx:02d} TINY     digit {txt!r} needs {gw:.2f}in, "
                    f"circle gives {inner_w:.2f}in")

        boxes.append((x + ml, y + mt, inner_w, need, txt, str(sh.shape_type)))

    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            ax, ay, aw, ah, at, _ = boxes[i]
            bx, by, bw, bh, bt, _ = boxes[j]
            ox = min(ax + aw, bx + bw) - max(ax, bx)
            oy = min(ay + ah, by + bh) - max(ay, by)
            if ox > 0.05 and oy > 0.05:
                findings.append(
                    f"s{idx:02d} OVERLAP  {ox:.2f}x{oy:.2f}in — "
                    f"{at[:26]!r} / {bt[:26]!r}")

print(f"{len(prs.slides)} slides, {len(findings)} findings\n")
for f in findings:
    print(f)
