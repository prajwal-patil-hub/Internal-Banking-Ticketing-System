"""Render the SOP deck to PDF, drawing directly from the .pptx shape tree.

Why this exists: a .pptx is a description that every viewer interprets for
itself, and weak mobile previewers interpret it badly — ghosting shapes
between slides, pairing a slide with the wrong picture. A PDF is not
interpreted, it is drawn. This produces the copy that looks the same
everywhere, and it is generated from the same file that is shipped, so the
two cannot drift apart.

Usage:
    python docs/sop/render-pdf.py <deck.pptx> <out.pdf>
"""
from __future__ import annotations

import io
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu
from reportlab.lib.colors import Color
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

SRC = Path(sys.argv[1])
OUT = Path(sys.argv[2])

EMU = 914400.0          # EMU per inch
PT = 72.0               # points per inch
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# Liberation Sans is metric-compatible with Arial and reads as a corporate
# sans. Calibri is not present here and cannot be embedded.
FONTS = {
    ("sans", False): ("SopSans", "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf"),
    ("sans", True): ("SopSans-Bold", "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf"),
}
for (_, bold), (name, path) in FONTS.items():
    pdfmetrics.registerFont(TTFont(name, path))


def fname(bold: bool) -> str:
    return FONTS[("sans", bool(bold))][0]


def col(c, default=None):
    try:
        if c is not None and c.type is not None and c.rgb is not None:
            r, g, b = tuple(c.rgb)
            return Color(r / 255, g / 255, b / 255)
    except Exception:
        pass
    return default


def geom_of(sh) -> str:
    try:
        g = sh._element.find(f".//{{{A}}}prstGeom")
        return g.get("prst") if g is not None else ""
    except Exception:
        return ""


def fill_of(sh):
    try:
        f = sh.fill
        if f.type == 1:
            return col(f.fore_color)
    except Exception:
        pass
    return None


def line_of(sh):
    try:
        ln = sh.line
        if ln.fill.type == 1:
            return col(ln.color)
    except Exception:
        pass
    return None


prs = Presentation(str(SRC))
SW = prs.slide_width / EMU * PT
SH = prs.slide_height / EMU * PT
c = canvas.Canvas(str(OUT), pagesize=(SW, SH))
c.setTitle("SUCCESS Bank — Internal Ticketing · Standard Operating Procedure")


def X(emu):
    return emu / EMU * PT


def Y(emu):
    """PDF origin is bottom-left; PowerPoint's is top-left."""
    return SH - emu / EMU * PT


def draw_shape_body(sh, x, y, w, h):
    g = geom_of(sh)
    f = fill_of(sh)
    ln = line_of(sh)
    if f is None and ln is None:
        return
    if f is not None:
        c.setFillColor(f)
    if ln is not None:
        c.setStrokeColor(ln)
        c.setLineWidth(0.75)
    mode = (1 if f is not None else 0, 1 if ln is not None else 0)

    if g == "ellipse":
        c.ellipse(x, y - h, x + w, y, stroke=mode[1], fill=mode[0])
    elif g == "roundRect":
        c.roundRect(x, y - h, w, h, min(w, h) * 0.16, stroke=mode[1], fill=mode[0])
    elif g in ("rightArrow", "leftArrow", "upArrow", "downArrow"):
        p = c.beginPath()
        if g == "rightArrow":
            pts = [(0, h / 3), (w * .65, h / 3), (w * .65, 0), (w, h / 2),
                   (w * .65, h), (w * .65, h * 2 / 3), (0, h * 2 / 3)]
        elif g == "leftArrow":
            pts = [(w, h / 3), (w * .35, h / 3), (w * .35, 0), (0, h / 2),
                   (w * .35, h), (w * .35, h * 2 / 3), (w, h * 2 / 3)]
        elif g == "upArrow":
            pts = [(w / 3, h), (w / 3, h * .35), (0, h * .35), (w / 2, 0),
                   (w, h * .35), (w * 2 / 3, h * .35), (w * 2 / 3, h)]
        else:
            pts = [(w / 3, 0), (w / 3, h * .65), (0, h * .65), (w / 2, h),
                   (w, h * .65), (w * 2 / 3, h * .65), (w * 2 / 3, 0)]
        p.moveTo(x + pts[0][0], y - pts[0][1])
        for px, py in pts[1:]:
            p.lineTo(x + px, y - py)
        p.close()
        c.drawPath(p, stroke=mode[1], fill=mode[0])
    else:
        c.rect(x, y - h, w, h, stroke=mode[1], fill=mode[0])


def wrap(text, font, size, avail):
    words, line, out = text.split(), "", []
    for wd in words:
        trial = (line + " " + wd).strip()
        if pdfmetrics.stringWidth(trial, font, size) <= avail or not line:
            line = trial
        else:
            out.append(line)
            line = wd
    out.append(line)
    return out


def draw_text_frame(tf, x, y, w, h, autoshape: bool):
    """Lay out a text frame the way PowerPoint would.

    Autoshapes centre their text vertically by default; text boxes do not.
    Getting that wrong is what puts a numeral off-centre in its circle.
    """
    def inset(v, default):
        return (v if v is not None else Emu(int(default * EMU))) / EMU * PT

    ml = inset(tf.margin_left, 0.1)
    mr = inset(tf.margin_right, 0.1)
    mt = inset(tf.margin_top, 0.05)
    mb = inset(tf.margin_bottom, 0.05)
    avail = max(1.0, w - ml - mr)

    anchor = str(tf.vertical_anchor or "")
    middle = anchor.startswith("MIDDLE") or (tf.vertical_anchor is None and autoshape)
    bottom = anchor.startswith("BOTTOM")

    lines = []          # (text, size, bold, colour, alignment)
    total = 0.0
    for p in tf.paragraphs:
        runs = [r for r in p.runs]
        if not runs:
            continue
        txt = "".join(r.text for r in runs)
        if not txt.strip():
            continue
        r0 = runs[0]
        size = r0.font.size.pt if r0.font.size else 18.0
        bold = bool(r0.font.bold)
        colour = col(r0.font.color, Color(0, 0, 0))
        align = str(p.alignment or "")
        sp = p.line_spacing if isinstance(p.line_spacing, float) else 1.15
        wrapped = wrap(txt, fname(bold), size, avail) if tf.word_wrap is not False else [txt]
        for w_line in wrapped:
            lines.append((w_line, size, bold, colour, align, size * sp))
            total += size * sp

    if not lines:
        return
    if middle:
        cur = y - mt - (h - mt - mb - total) / 2
    elif bottom:
        cur = y - h + mb + total
    else:
        cur = y - mt

    for text, size, bold, colour, align, adv in lines:
        f = fname(bold)
        tw = pdfmetrics.stringWidth(text, f, size)
        if align.startswith("CENTER"):
            tx = x + ml + (avail - tw) / 2
        elif align.startswith("RIGHT"):
            tx = x + w - mr - tw
        else:
            tx = x + ml
        c.setFont(f, size)
        c.setFillColor(colour)
        # Baseline sits ~0.8 of the em below the line top, which matches how
        # PowerPoint seats a single line inside its box closely enough that a
        # numeral lands in the middle of its circle.
        c.drawString(tx, cur - size * 0.80, text)
        cur -= adv


def draw_table(sh, x, y):
    t = sh.table
    ry = y
    for r_i, row in enumerate(t.rows):
        rh = row.height / EMU * PT
        cx = x
        for c_i, cell in enumerate(row.cells):
            cw = t.columns[c_i].width / EMU * PT
            fc = None
            try:
                fc = col(cell.fill.fore_color, Color(1, 1, 1))
            except Exception:
                fc = Color(1, 1, 1)
            c.setFillColor(fc)
            c.setStrokeColor(Color(0.84, 0.82, 0.76))
            c.setLineWidth(0.5)
            c.rect(cx, ry - rh, cw, rh, stroke=1, fill=1)
            p = cell.text_frame.paragraphs[0]
            if p.runs:
                r0 = p.runs[0]
                size = r0.font.size.pt if r0.font.size else 11.0
                f = fname(bool(r0.font.bold))
                c.setFont(f, size)
                c.setFillColor(col(r0.font.color, Color(0.1, 0.16, 0.2)))
                txt = cell.text_frame.text
                for ln in wrap(txt, f, size, cw - 10):
                    c.drawString(cx + 5, ry - size * 1.25, ln)
                    ry -= 0  # single line per cell in this deck
            cx += cw
        ry -= rh


def slide_bg(slide):
    """The slide's own background colour, if it sets one."""
    try:
        f = slide.background.fill
        if f.type == 1:
            return col(f.fore_color)
    except Exception:
        pass
    return None


for idx, slide in enumerate(prs.slides, start=1):
    bg = slide_bg(slide) or Color(1, 1, 1)
    c.setFillColor(bg)
    c.rect(0, 0, SW, SH, stroke=0, fill=1)

    for sh in slide.shapes:
        if sh.left is None or sh.width is None:
            continue
        x, y = X(sh.left), Y(sh.top)
        w, h = X(sh.width), sh.height / EMU * PT

        if sh.shape_type == 13:                      # picture
            try:
                im = Image.open(io.BytesIO(sh.image.blob)).convert("RGB")
                c.drawImage(ImageReader(im), x, y - h, w, h)
            except Exception:
                pass
            continue

        if sh.has_table:
            draw_table(sh, x, y)
            continue

        autoshape = str(sh.shape_type).startswith("AUTO_SHAPE")
        if autoshape:
            draw_shape_body(sh, x, y, w, h)

        if sh.has_text_frame and sh.text_frame.text.strip():
            draw_text_frame(sh.text_frame, x, y, w, h, autoshape)

    c.showPage()

c.save()
print(f"wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} pages)")
