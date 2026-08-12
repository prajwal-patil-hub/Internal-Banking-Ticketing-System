"""Render slides from a .pptx to PNG, reading the file's own shape tree.

Not a PowerPoint-accurate renderer — it is a proof sheet. It draws what is
actually stored in the file (positions, sizes, fills, text, images, and the
text insets that PowerPoint will honour), which is what catches a numeral
squeezed out of its circle or a label overrunning its box.
"""
from __future__ import annotations

import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu

SRC = Path(sys.argv[1])
OUT = Path(sys.argv[2])
ONLY = [int(x) for x in sys.argv[3].split(",")] if len(sys.argv) > 3 else None
OUT.mkdir(parents=True, exist_ok=True)

SCALE = 110 / 914400  # px per EMU  → ~1466 px wide slide

FONTS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
]
BOLDS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
]


def _font(size_pt: float, bold: bool):
    px = max(7, int(size_pt * SCALE * 12700))
    for p in (BOLDS if bold else FONTS):
        if Path(p).exists():
            try:
                return ImageFont.truetype(p, px)
            except Exception:
                pass
    return ImageFont.load_default()


def rgb(c, default=(120, 120, 120)):
    try:
        if c and c.type is not None and c.rgb is not None:
            return tuple(c.rgb)
    except Exception:
        pass
    return default


prs = Presentation(str(SRC))
SW, SH = int(prs.slide_width * SCALE), int(prs.slide_height * SCALE)

for idx, slide in enumerate(prs.slides, start=1):
    if ONLY and idx not in ONLY:
        continue
    img = Image.new("RGB", (SW, SH), (255, 255, 255))
    d = ImageDraw.Draw(img)

    for sh in slide.shapes:
        if sh.left is None:
            continue
        x, y = int(sh.left * SCALE), int(sh.top * SCALE)
        w, h = int((sh.width or 0) * SCALE), int((sh.height or 0) * SCALE)

        # pictures
        if sh.shape_type == 13:
            try:
                im = Image.open(__import__("io").BytesIO(sh.image.blob)).convert("RGB")
                img.paste(im.resize((max(1, w), max(1, h))), (x, y))
            except Exception:
                d.rectangle([x, y, x + w, y + h], outline=(200, 0, 0), width=2)
            continue

        # fills
        is_oval = "OVAL" in str(sh.shape_type) or (
            sh.shape_type == 1 and getattr(getattr(sh, "adjustments", None), "_adjustments", None) is not None
            and w == h and w < int(0.5 * 914400 * SCALE))
        try:
            fill = sh.fill
            col = rgb(fill.fore_color, None) if fill.type == 1 else None
        except Exception:
            col = None
        if col:
            shape_name = str(sh.shape_type)
            try:
                st = sh._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom")
                geom = st.get("prst") if st is not None else ""
            except Exception:
                geom = ""
            if geom == "ellipse":
                d.ellipse([x, y, x + w, y + h], fill=col)
            elif geom == "roundRect":
                d.rounded_rectangle([x, y, x + w, y + h], radius=max(3, h // 6), fill=col)
            elif geom == "rightArrow":
                d.polygon([(x, y + h//3), (x + int(w*0.65), y + h//3), (x + int(w*0.65), y),
                           (x + w, y + h//2), (x + int(w*0.65), y + h),
                           (x + int(w*0.65), y + int(h*0.66)), (x, y + int(h*0.66))], fill=col)
            elif geom == "leftArrow":
                d.polygon([(x + w, y + h//3), (x + int(w*0.35), y + h//3), (x + int(w*0.35), y),
                           (x, y + h//2), (x + int(w*0.35), y + h),
                           (x + int(w*0.35), y + int(h*0.66)), (x + w, y + int(h*0.66))], fill=col)
            elif geom == "upArrow":
                d.polygon([(x + w//3, y + h), (x + w//3, y + int(h*0.35)), (x, y + int(h*0.35)),
                           (x + w//2, y), (x + w, y + int(h*0.35)),
                           (x + int(w*0.66), y + int(h*0.35)), (x + int(w*0.66), y + h)], fill=col)
            elif geom == "downArrow":
                d.polygon([(x + w//3, y), (x + w//3, y + int(h*0.65)), (x, y + int(h*0.65)),
                           (x + w//2, y + h), (x + w, y + int(h*0.65)),
                           (x + int(w*0.66), y + int(h*0.65)), (x + int(w*0.66), y)], fill=col)
            else:
                d.rectangle([x, y, x + w, y + h], fill=col)

        # table
        if sh.has_table:
            t = sh.table
            ry = y
            for r_i, row in enumerate(t.rows):
                rh = int(row.height * SCALE)
                cx = x
                for c_i, cell in enumerate(row.cells):
                    cw = int(t.columns[c_i].width * SCALE)
                    try:
                        cc = rgb(cell.fill.fore_color, (255, 255, 255))
                    except Exception:
                        cc = (255, 255, 255)
                    d.rectangle([cx, ry, cx + cw, ry + rh], fill=cc, outline=(215, 208, 195))
                    p = cell.text_frame.paragraphs[0]
                    if p.runs:
                        r0 = p.runs[0]
                        f = _font(r0.font.size.pt if r0.font.size else 11, bool(r0.font.bold))
                        d.text((cx + 6, ry + 5), cell.text_frame.text[:70],
                               font=f, fill=rgb(r0.font.color, (27, 42, 51)))
                    cx += cw
                ry += rh
            continue

        # text
        if not sh.has_text_frame:
            continue
        tf = sh.text_frame
        if not tf.text.strip():
            continue
        ml = int((tf.margin_left or 0) * SCALE)
        mr = int((tf.margin_right or 0) * SCALE)
        mt = int((tf.margin_top or 0) * SCALE)
        anchor_mid = str(tf.vertical_anchor) .startswith("MIDDLE")

        ty = y + mt
        if anchor_mid:
            total = 0
            for p in tf.paragraphs:
                sz = p.runs[0].font.size.pt if p.runs and p.runs[0].font.size else 12
                total += int(sz * SCALE * 12700 * 1.25)
            ty = y + max(0, (h - total) // 2)

        for p in tf.paragraphs:
            if not p.runs:
                ty += 10
                continue
            r0 = p.runs[0]
            sz = r0.font.size.pt if r0.font.size else 12
            f = _font(sz, bool(r0.font.bold))
            txt = "".join(r.text for r in p.runs)
            col = rgb(r0.font.color, (27, 42, 51))
            avail = max(4, w - ml - mr)
            # wrap
            words, line, lines = txt.split(" "), "", []
            for wd in words:
                trial = (line + " " + wd).strip()
                if d.textlength(trial, font=f) <= avail or not line:
                    line = trial
                else:
                    lines.append(line); line = wd
            lines.append(line)
            for ln in lines:
                tw = d.textlength(ln, font=f)
                if str(p.alignment).startswith("CENTER"):
                    tx = x + ml + (avail - tw) / 2
                elif str(p.alignment).startswith("RIGHT"):
                    tx = x + w - mr - tw
                else:
                    tx = x + ml
                d.text((tx, ty), ln, font=f, fill=col)
                ty += int(sz * SCALE * 12700 * 1.25)

    img.save(OUT / f"slide-{idx:02d}.png")

print(f"rendered -> {OUT}")
