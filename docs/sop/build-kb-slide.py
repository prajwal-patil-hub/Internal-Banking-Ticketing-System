"""Build the Knowledge Base slide as a standalone deck.

Same visual language as the SOP deck so it drops straight in, but
self-contained: one slide, its own file.

The screenshot is cropped into two panels from a single real capture — the
answering half on top, the curating half below — because the page is taller
than it is wide and scaling the whole thing to slide height would render the
type unreadably small. Both panels are pixels from the same render; nothing
here is a mock-up.

    python docs/sop/build-kb-slide.py <shots-dir> <out.pptx>
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

SHOTS = Path(sys.argv[1] if len(sys.argv) > 1 else "docs/sop/screens")
OUT = Path(sys.argv[2] if len(sys.argv) > 2 else "docs/sop/Knowledge-Base-slide.pptx")

TEAL = RGBColor(0x0E, 0x4F, 0x4A)
CREAM = RGBColor(0xF3, 0xEC, 0xE0)
INK = RGBColor(0x1B, 0x2A, 0x33)
MUTE = RGBColor(0x5B, 0x6B, 0x75)
LINE = RGBColor(0xD6, 0xCE, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xC2, 0x5B, 0x2E)
OK = RGBColor(0x1E, 0x7A, 0x53)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
s = prs.slides.add_slide(prs.slide_layouts[6])
bg = s.background.fill
bg.solid()
bg.fore_color.rgb = WHITE


def rect(x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sh = s.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def tb(x, y, w, h, text, size=12, bold=False, color=INK,
       align=PP_ALIGN.LEFT, spacing=1.2, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    # Every paragraph is styled, not just the first: `tf.text` with newlines
    # creates separate paragraphs and an unstyled one falls back to 18pt black.
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.line_spacing = spacing
        for r in p.runs:
            r.font.size, r.font.bold, r.font.color.rgb, r.font.name = (
                Pt(size), bold, color, "Calibri")
    return box


def marker(cx, cy, n, d, fill, fontsize):
    """Numbered circle with the digit actually centred — zero insets, no autofit."""
    sh = rect(cx - d / 2, cy - d / 2, d, d, fill, shape=MSO_SHAPE.OVAL)
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.text = str(n)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.0
    r = p.runs[0]
    r.font.size, r.font.bold, r.font.color.rgb, r.font.name = (
        Pt(fontsize), True, WHITE, "Calibri")


def place(src: Image.Image, box: tuple[int, int, int, int], x, y, w):
    """Crop, downsample to 200dpi and place. Returns the drawn height."""
    crop = src.crop(box)
    target_px = int(w / 914400 * 200)
    im = crop.convert("RGB").resize(
        (target_px, round(crop.height * target_px / crop.width)), Image.LANCZOS
    )
    im = im.quantize(colors=256, dither=Image.FLOYDSTEINBERG)
    buf = io.BytesIO()
    im.save(buf, format="PNG", optimize=True)
    buf.seek(0)
    h = Emu(int(w * crop.height / crop.width))
    rect(x - Inches(0.022), y - Inches(0.022), w + Inches(0.044), h + Inches(0.044), LINE)
    s.shapes.add_picture(buf, x, y, width=w)
    return h


# ---- header ----------------------------------------------------------------
rect(0, 0, W, Inches(0.95), TEAL)
rect(0, Inches(0.95), W, Inches(0.04), ACCENT)
tb(Inches(0.55), Inches(0.17), Inches(6), Inches(0.3),
   "THE KNOWLEDGE BASE", size=12, bold=True, color=ACCENT)
# Width stops short of the role label's box at x=9.03. The text is shorter
# than that today, so the two never touched visually — but overlapping boxes
# are a collision waiting for a longer title, which is what the layout audit
# flags and why it is worth fixing while it is still cosmetic.
tb(Inches(0.55), Inches(0.46), Inches(8.30), Inches(0.4),
   "Asking it a question, and keeping it stocked", size=20, bold=True, color=WHITE)
tb(W - Inches(4.3), Inches(0.35), Inches(3.75), Inches(0.3),
   "ASK: AGENT & ABOVE  ·  UPLOAD: ADMIN", size=11, bold=True, color=CREAM,
   align=PP_ALIGN.RIGHT)

# ---- panels ----------------------------------------------------------------
shot = Image.open(SHOTS / "70-knowledge-base.png")
SW, SH = shot.size

# Markers for left-aligned targets sit in a gutter to the left of the panel.
# Drawing them over the panel put them on top of the very labels they point at
# — a marker that obscures its subject is worse than no marker.
PX, PW = Inches(0.80), Inches(6.75)
GUTTER_CX = Inches(0.60)

# Panel A — the answering half: question, confidence, answer, sources.
# Cropped past the sidebar (x) and past the stat strip (y), so the type stays
# legible at slide scale instead of being shrunk to fit decoration.
A_BOX = (int(SW * 0.178), int(SH * 0.221), SW, int(SH * 0.590))
A_Y = Inches(1.22)
a_h = place(shot, A_BOX, PX, A_Y, PW)

# Panel B — the curating half: collections, grants, document states.
B_BOX = (int(SW * 0.178), int(SH * 0.609), SW, int(SH * 0.960))
B_Y = A_Y + a_h + Inches(0.13)
b_h = place(shot, B_BOX, PX, B_Y, PW)

# (panel, x-fraction or None for the gutter, y-fraction, label). Reading order
# is top to bottom within a panel, and the answering panel comes first because
# that is the order a reader meets the feature: you ask before you curate.
CALLOUTS = [
    ("A", None,  0.204, "Ask in plain words — no query syntax"),
    ("A", None,  0.312, "Confidence band, computed from the evidence\nrather than claimed by the model"),
    ("A", None,  0.428, "Every factual sentence carries its source [n]"),
    ("A", None,  0.522, "Cited sources: document, section and page"),
    ("A", None,  0.852, "Also retrieved but not cited — shown so you\ncan see what was weighed"),
    ("B", 0.700, 0.087, "Admins upload PDF, Word, Markdown, text or CSV"),
    ("B", 0.500, 0.307, "Which roles may search this collection"),
    ("B", None,  0.324, "A collection nobody can read is flagged,\nnot left to fail silently"),
    ("B", None,  0.907, "A document that failed to index says why"),
]

D = Inches(0.28)
for i, (panel, fx, fy, _label) in enumerate(CALLOUTS, start=1):
    py, ph = (A_Y, a_h) if panel == "A" else (B_Y, b_h)
    cy = py + Emu(int(ph * fy))
    cx = GUTTER_CX if fx is None else PX + Emu(int(PW * fx))
    marker(cx, cy, i, D, ACCENT, 12)

# ---- right column ----------------------------------------------------------
rx, rw = Inches(7.85), Inches(5.03)
y = Inches(1.22)

tb(rx, y, rw, Inches(0.24), "ON THIS SCREEN", size=10, bold=True, color=ACCENT)
y += Inches(0.30)
for i, (_p, _fx, _fy, label) in enumerate(CALLOUTS, start=1):
    lines = label.count("\n") + 1
    h = Inches(0.175) * lines
    marker(rx + Inches(0.105), y + Inches(0.088), i, Inches(0.21), ACCENT, 9)
    tb(rx + Inches(0.33), y, rw - Inches(0.33), h, label, size=10, color=INK, spacing=1.05)
    y += h + Inches(0.075)

y += Inches(0.14)
rect(rx, y, Inches(0.045), Inches(1.30), OK)
tb(rx + Inches(0.17), y, rw - Inches(0.17), Inches(0.18),
   "WHY YOU CAN TRUST THE ANSWER", size=9.5, bold=True, color=OK)
tb(rx + Inches(0.17), y + Inches(0.22), rw - Inches(0.17), Inches(1.05),
   "Retrieval is filtered by your role before the model sees anything, so it "
   "cannot quote a document you could not open yourself. Every citation is "
   "checked against what was actually retrieved — an invented one is removed "
   "along with the sentence it supported.",
   size=10, color=INK, spacing=1.12)
y += Inches(1.46)

rect(rx, y, Inches(0.045), Inches(0.82), ACCENT)
tb(rx + Inches(0.17), y, rw - Inches(0.17), Inches(0.18),
   "WHEN IT HAS NOTHING", size=9.5, bold=True, color=ACCENT)
tb(rx + Inches(0.17), y + Inches(0.22), rw - Inches(0.17), Inches(0.58),
   "It says so and stops, rather than guessing. That is the correct answer, "
   "not a failure.",
   size=10, color=INK, spacing=1.12)

# ---- footer ----------------------------------------------------------------
rect(0, H - Inches(0.42), W, Inches(0.42), CREAM)
tb(Inches(0.55), H - Inches(0.34), Inches(9), Inches(0.25),
   "SUCCESS Bank — Internal Ticketing · Standard Operating Procedure",
   size=9, color=MUTE)
tb(W - Inches(2.4), H - Inches(0.34), Inches(1.85), Inches(0.25),
   "Knowledge Base", size=9, color=MUTE, align=PP_ALIGN.RIGHT)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"saved {OUT} (1 slide, {len(CALLOUTS)} callouts)")
