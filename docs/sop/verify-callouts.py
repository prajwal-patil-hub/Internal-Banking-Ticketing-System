"""Check callout markers and shape placement — geometry the text audit cannot see.

`audit-layout.py` asks whether text fits its box. This asks whether things are
in the right *place*: markers inside the screenshot they annotate, not on top
of each other, numbered in reading order, and no shape straying into the
footer or off the slide.

Reported as three severities:

  ERROR   wrong: a marker outside its screenshot, a shape off the slide
  WARN    likely wrong: markers overlapping, numbering out of reading order
  INFO    worth a look: a marker very close to another
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

SRC = Path(sys.argv[1])
EMU = 914400.0
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

#: The footer band occupies the bottom 0.42in of every slide.
FOOTER_H = 0.42

findings: list[tuple[str, int, str]] = []


def add(sev: str, slide: int, msg: str) -> None:
    findings.append((sev, slide, msg))


def geom(sh) -> str:
    try:
        g = sh._element.find(f".//{{{A}}}prstGeom")
        return g.get("prst") if g is not None else ""
    except Exception:
        return ""


def box(sh):
    return (sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU)


def overlap(a, b) -> float:
    """Intersection area of two (x, y, w, h) boxes, in square inches."""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ox = min(ax + aw, bx + bw) - max(ax, bx)
    oy = min(ay + ah, by + bh) - max(ay, by)
    return ox * oy if ox > 0 and oy > 0 else 0.0


prs = Presentation(str(SRC))
SW, SH = prs.slide_width / EMU, prs.slide_height / EMU
footer_top = SH - FOOTER_H

for idx, slide in enumerate(prs.slides, start=1):
    pictures, markers, others = [], [], []

    for sh in slide.shapes:
        if sh.left is None or sh.width is None:
            continue
        b = box(sh)

        if sh.left < -0.01 or sh.top < -0.01 or b[0] + b[2] > SW + 0.01 or b[1] + b[3] > SH + 0.01:
            add("ERROR", idx, f"shape extends off the slide: {b[0]:.2f},{b[1]:.2f} {b[2]:.2f}x{b[3]:.2f}")

        if sh.shape_type == 13:
            pictures.append((sh, b))
            continue

        txt = sh.text_frame.text.strip() if sh.has_text_frame else ""
        if geom(sh) == "ellipse" and txt.isdigit() and b[2] < 0.5:
            markers.append((int(txt), b))
        else:
            others.append((sh, b, txt))

    # --- markers must sit on the screenshot they annotate --------------------
    # Legend markers live in the right-hand column and are excluded by x.
    img = pictures[0][1] if pictures else None
    on_image = [(n, b) for n, b in markers if img and b[0] < img[0] + img[2]]

    for n, b in on_image:
        cx, cy = b[0] + b[2] / 2, b[1] + b[3] / 2
        inside = (img[0] - 0.2 <= cx <= img[0] + img[2] + 0.05
                  and img[1] - 0.05 <= cy <= img[1] + img[3] + 0.05)
        if not inside:
            add("ERROR", idx,
                f"marker {n} at ({cx:.2f},{cy:.2f}) is outside the screenshot "
                f"({img[0]:.2f},{img[1]:.2f} {img[2]:.2f}x{img[3]:.2f})")

    # --- markers must not sit on top of one another -------------------------
    for i in range(len(on_image)):
        for j in range(i + 1, len(on_image)):
            n1, b1 = on_image[i]
            n2, b2 = on_image[j]
            ov = overlap(b1, b2)
            area = b1[2] * b1[3]
            if ov > area * 0.25:
                add("WARN", idx, f"markers {n1} and {n2} overlap by {ov / area:.0%}")
            elif ov > 0:
                add("INFO", idx, f"markers {n1} and {n2} touch")

    # --- numbering must follow reading order --------------------------------
    # Rows are clustered exactly as the builder clusters them. Checking with a
    # fixed grid while the builder clusters would be two implementations of
    # one rule, which is the defect this file exists to catch.
    if len(on_image) > 1:
        ROW_GAP = 0.55                      # inches on the slide
        rows: list[list] = []
        for m in sorted(on_image, key=lambda m: m[1][1] + m[1][3] / 2):
            cy = m[1][1] + m[1][3] / 2
            if rows and cy - (rows[-1][0][1][1] + rows[-1][0][1][3] / 2) <= ROW_GAP:
                rows[-1].append(m)
            else:
                rows.append([m])
        ordered = [m for row in rows for m in sorted(row, key=lambda m: m[1][0])]
        if [n for n, _ in ordered] != sorted(n for n, _ in on_image):
            add("WARN", idx,
                "markers are not numbered in reading order: reading order is "
                + ", ".join(str(n) for n, _ in ordered))

    # --- nothing may stray into the footer band -----------------------------
    for sh, b, txt in others:
        if b[3] >= FOOTER_H * 0.9 and abs(b[1] - footer_top) < 0.02:
            continue                      # the footer bar itself
        if b[1] <= 0.01 and b[1] + b[3] >= SH - 0.01:
            continue                      # a full-height decorative spine
        if b[1] + b[3] > footer_top + 0.02 and b[1] < footer_top:
            add("ERROR", idx,
                f"shape crosses into the footer band (bottom {b[1] + b[3]:.2f}in "
                f"vs footer at {footer_top:.2f}in): {txt[:40]!r}")

    for sh, b in pictures:
        if b[1] + b[3] > footer_top + 0.02:
            add("ERROR", idx, f"screenshot runs into the footer (bottom {b[1] + b[3]:.2f}in)")


order = {"ERROR": 0, "WARN": 1, "INFO": 2}
findings.sort(key=lambda f: (order[f[0]], f[1]))
errs = sum(1 for f in findings if f[0] == "ERROR")
warns = sum(1 for f in findings if f[0] == "WARN")

print(f"{len(prs.slides._sldIdLst)} slides — {errs} errors, {warns} warnings, "
      f"{len(findings) - errs - warns} notes\n")
for sev, slide, msg in findings:
    print(f"  {sev:5} s{slide:02d}  {msg}")
if not findings:
    print("  nothing to report")
